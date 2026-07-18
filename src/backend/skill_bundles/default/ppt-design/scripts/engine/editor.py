"""PowerPoint slide-level editing primitives (python-pptx).

Functions are designed to be invoked one at a time from the MCP layer; each
loads the deck, performs ONE op, saves, and reports. State is the file_id
returned to the LLM — there is no in-memory deck cache between calls.
"""
from __future__ import annotations

from typing import Any

from ._shims import input_path, output_path

# 16:9 default canvas (matches builder.build_from_spec)
_INCH = 914400
_DEFAULT_W = 10 * _INCH
_DEFAULT_H = int(5.625 * _INCH)


def _check_slide_index(prs, slide_index: int) -> None:
    n = len(prs.slides)
    if not (0 <= slide_index < n):
        raise ValueError(
            f"slide_index {slide_index} out of range (presentation has {n} slides)"
        )


def _hex_to_rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str.lstrip("#"))


def set_slide_title(
    *,
    input_filename: str,
    output_filename: str,
    slide_index: int,
    title: str,
) -> dict[str, Any]:
    """Set or replace the title text on an existing slide.

    Strategy: locate the first text frame whose name contains "title" or whose
    placeholder type is TITLE; if none found, insert a new title text box at
    the canonical 16:9 title position.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(input_path(input_filename)))
    _check_slide_index(prs, slide_index)
    slide = prs.slides[slide_index]

    title_shape = None
    # 1. Try placeholder "title" / "centerTitle"
    for shape in slide.shapes:
        if shape.has_text_frame and (
            getattr(shape, "is_placeholder", False)
            and shape.placeholder_format
            and shape.placeholder_format.idx == 0
        ):
            title_shape = shape
            break

    if title_shape is None:
        # 2. Try first text frame at top of slide (heuristic: y < 1.2 inches)
        for shape in slide.shapes:
            if shape.has_text_frame and shape.top is not None and shape.top < int(1.2 * _INCH):
                title_shape = shape
                break

    if title_shape is None:
        # 3. Insert a new text box at the canonical title slot
        box = slide.shapes.add_textbox(
            Emu(int(0.5 * _INCH)),
            Emu(int(0.4 * _INCH)),
            Emu(int(9.0 * _INCH)),
            Emu(int(0.8 * _INCH)),
        )
        title_shape = box

    tf = title_shape.text_frame
    tf.text = title

    out = output_path(output_filename)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    return {
        "output_filename": output_filename,
        "slide_index": slide_index,
        "title": title,
    }


def add_text_to_slide(
    *,
    input_filename: str,
    output_filename: str,
    slide_index: int,
    text: str,
    position: dict | None = None,
    font_size: int = 14,
    color_hex: str | None = None,
    bold: bool = False,
) -> dict[str, Any]:
    """Insert a free-form text box onto a slide.

    Args:
        position: optional ``{"left", "top", "width", "height"}`` in inches;
            default places the box mid-slide. Floats accepted.
        font_size: default 14 pt
        color_hex: optional 6-char hex (no leading ``#``)
        bold: bold the inserted text
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt

    prs = Presentation(str(input_path(input_filename)))
    _check_slide_index(prs, slide_index)
    slide = prs.slides[slide_index]

    pos = position or {}
    left = Emu(int(float(pos.get("left", 1.0)) * _INCH))
    top = Emu(int(float(pos.get("top", 2.5)) * _INCH))
    width = Emu(int(float(pos.get("width", 8.0)) * _INCH))
    height = Emu(int(float(pos.get("height", 1.0)) * _INCH))

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text

    para = tf.paragraphs[0]
    para.font.size = Pt(int(font_size))
    para.font.bold = bool(bold)
    if color_hex:
        para.font.color.rgb = _hex_to_rgb(color_hex)
    for run in para.runs:
        run.font.size = Pt(int(font_size))
        run.font.bold = bool(bold)
        if color_hex:
            run.font.color.rgb = _hex_to_rgb(color_hex)

    out = output_path(output_filename)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    return {
        "output_filename": output_filename,
        "slide_index": slide_index,
        "text_chars": len(text),
        "position": {
            "left": float(pos.get("left", 1.0)),
            "top": float(pos.get("top", 2.5)),
            "width": float(pos.get("width", 8.0)),
            "height": float(pos.get("height", 1.0)),
        },
    }


def insert_image(
    *,
    input_filename: str,
    output_filename: str,
    slide_index: int,
    image_filename: str,
    position: dict | None = None,
) -> dict[str, Any]:
    """Insert an image onto a slide.

    Args:
        image_filename: filename in cwd (the MCP tool ships the image bytes via
            ``input_files_b64`` so they appear in cwd alongside the .pptx).
        position: ``{"left", "top", "width"?, "height"?}`` in inches; if width
            and height omitted, the image is inserted at native size.
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(input_path(input_filename)))
    _check_slide_index(prs, slide_index)
    slide = prs.slides[slide_index]

    pos = position or {}
    left = Emu(int(float(pos.get("left", 1.0)) * _INCH))
    top = Emu(int(float(pos.get("top", 1.0)) * _INCH))
    width = Emu(int(float(pos["width"]) * _INCH)) if "width" in pos else None
    height = Emu(int(float(pos["height"]) * _INCH)) if "height" in pos else None

    img_path = str(input_path(image_filename))
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    out = output_path(output_filename)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    return {
        "output_filename": output_filename,
        "slide_index": slide_index,
        "image_filename": image_filename,
        "rendered_width_emu": pic.width,
        "rendered_height_emu": pic.height,
    }


def delete_slide(
    *,
    input_filename: str,
    output_filename: str,
    slide_index: int,
) -> dict[str, Any]:
    """Delete a slide by 0-based index.

    python-pptx has no built-in ``delete_slide`` — we manipulate the underlying
    XML to drop the slide id and the relationship. This is the standard pattern
    documented in python-pptx issues (e.g. scanny/python-pptx#67).
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn

    prs = Presentation(str(input_path(input_filename)))
    _check_slide_index(prs, slide_index)

    # XML drop: 1) remove sldId from <p:sldIdLst>, 2) drop the relationship
    sldIdLst = prs.slides._sldIdLst  # internal accessor
    sld_ids = list(sldIdLst)
    target = sld_ids[slide_index]
    rId = target.attrib[qn("r:id")]
    sldIdLst.remove(target)
    prs.part.drop_rel(rId)

    out = output_path(output_filename)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    return {
        "output_filename": output_filename,
        "deleted_slide_index": slide_index,
        "remaining_slide_count": len(prs.slides),
    }
