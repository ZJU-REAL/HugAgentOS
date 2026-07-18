"""QA helpers for generated .pptx files.

Pipeline that mirrors the user's pptx skill QA loop:
  1. ``extract_all_text(path)`` — flatten every slide's text for review
  2. ``check_placeholders(path)`` — grep for "xxxx", "lorem", "占位", etc.

Called from the ``extract`` and ``check-placeholders`` subcommands of the
skill's CLI so the agent can run a fix-and-verify cycle after building a deck.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from ._shims import input_path


DEFAULT_PLACEHOLDER_PATTERNS: list[str] = [
    "xxxx",
    "lorem",
    "ipsum",
    "placeholder",
    "TODO",
    "占位",
    "待补充",
    "请填写",
    "this page layout",
    "this slide layout",
]

# Precompiled at module import — DEFAULT_PLACEHOLDER_PATTERNS is constant, no
# need to recompile on every check_placeholders call.
_DEFAULT_COMPILED: list[tuple[str, "re.Pattern[str]"]] = [
    (p, re.compile(re.escape(p), re.IGNORECASE)) for p in DEFAULT_PLACEHOLDER_PATTERNS
]


def _iter_text_runs(slide) -> list[str]:
    out: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            tf = shape.text_frame
        except Exception:
            continue
        for para in tf.paragraphs:
            text = "".join(run.text for run in para.runs) or para.text or ""
            text = text.strip()
            if text:
                out.append(text)
    return out


def extract_all_text(file_path: str | Path) -> list[dict[str, Any]]:
    """Return ``[{slide_index (1-based), text}]`` for every text-bearing paragraph."""
    from pptx import Presentation

    p = Path(file_path) if Path(file_path).is_absolute() else input_path(str(file_path))
    prs = Presentation(str(p))

    out: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        for text in _iter_text_runs(slide):
            out.append({"slide_index": idx + 1, "text": text})
    return out


def check_placeholders(
    file_path: str | Path,
    *,
    patterns: list[str] | None = None,
    text_blocks: list[dict[str, Any]] | None = None,
) -> list[dict[str, Any]]:
    """Scan slide text for placeholder strings.

    ``text_blocks`` lets callers reuse a previous ``extract_all_text`` result
    and avoid re-parsing the OOXML — useful when both QA tools run in sequence.
    """
    if patterns is None:
        compiled = _DEFAULT_COMPILED
    elif not patterns:
        return []
    else:
        compiled = [(p, re.compile(re.escape(p), re.IGNORECASE)) for p in patterns]

    blocks = text_blocks if text_blocks is not None else extract_all_text(file_path)

    hits: list[dict[str, Any]] = []
    for entry in blocks:
        text = entry["text"]
        for label, regex in compiled:
            if regex.search(text):
                hits.append({
                    "slide_index": entry["slide_index"],
                    "text": text,
                    "pattern": label,
                })
                break
    return hits
