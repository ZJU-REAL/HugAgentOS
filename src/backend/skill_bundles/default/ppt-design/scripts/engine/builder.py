"""PowerPoint creation + slide manipulation.

Two engines are supported:

  - ``pptxgenjs``   — Node.js-based canonical engine (default for ``build_from_spec``);
                       full visual fidelity, all packs + ~70 components
  - ``python-pptx`` — reduced pure-Python renderer: powers single-slide editing
                       (``add_slide``) and a no-Node-runtime build backup.
                       Does NOT render the rich WP16/WP17 components (degrades
                       to bullets) — the JS engine is the source of truth.

``build_from_spec`` is the primary entry point. ``add_slide`` is preserved
for editing existing decks (single-slide append/delete via the python-pptx
renderer).
"""
from __future__ import annotations

from typing import Any

from ._shims import input_path, output_path
from . import slide_types
from .style_recipes import get_recipe, list_recipes
from .themes import get_palette


VALID_SLIDE_TYPES = ("cover", "toc", "section", "content", "summary", "closing", "thanks")


def _new_presentation():
    """Create a 16:9 (10×5.625 in) python-pptx Presentation."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs


def _blank_layout(prs):
    """Return a safe blank slide layout.

    ``slide_layouts[6]`` is "blank" in python-pptx's default template, but a
    deck built by the pptxgenjs engine (the default `build` engine) carries a
    different master with **fewer** layouts — so the hard-coded index 6 raised
    ``IndexError`` when ``add-slide`` ran on a normally-generated deck. Pick the
    layout with the fewest placeholders (the real "blank" one), falling back to
    index 6 then the last layout.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("presentation has no slide layouts")
    blank = min(layouts, key=lambda lo: len(lo.placeholders))
    if len(blank.placeholders) == 0:
        return blank
    return layouts[6] if len(layouts) > 6 else layouts[-1]


# ── Public: spec-based builder (primary entry point) ─────────────────


def build_from_spec(
    *,
    spec: dict[str, Any],
    output_filename: str,
    engine: str = "pptxgenjs",
    style: str = "soft",
    theme: str | None = None,
) -> dict[str, Any]:
    """Build a complete .pptx from a JSON-like spec.

    Args:
        spec: presentation spec — same schema as
            ``scripts/engine/node_scripts/build_presentation.js``::

              {
                "title"?, "author"?, "subject"?,
                "theme"?: <name string> | { primary, secondary, accent, light, bg },
                "slides": [
                  {"type": "cover",   "title", "subtitle"?, "body"?},
                  {"type": "toc",     "title"?, "items": [...]},
                  {"type": "section", "title", "subtitle"?},
                  {"type": "content", "title", "bullets"? | leftBullets+rightBullets,
                                      "highlights"? <=3, "body"?},
                  {"type": "summary", "title"?, "bullets"?, "body"?}
                ]
              }

        output_filename: relative output filename (resolved via the engine's workdir)
        engine: ``"pptxgenjs"`` (default) or ``"python-pptx"``
        style: visual style — ``sharp`` / ``soft`` / ``rounded`` / ``pill``
        theme: overrides ``spec.theme`` if given (palette name or hex dict)

    Returns:
        ``{output_filename, engine, theme, style, slide_count, size_bytes, ...}``
    """
    if not isinstance(spec, dict):
        raise ValueError("spec must be a dict")

    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("spec.slides must be a non-empty list")

    for idx, sl in enumerate(slides):
        if not isinstance(sl, dict):
            raise ValueError(f"spec.slides[{idx}] must be a dict")
        st = str(sl.get("type") or "content").lower()
        if st not in VALID_SLIDE_TYPES:
            raise ValueError(
                f"spec.slides[{idx}].type={st!r} invalid; expected one of {VALID_SLIDE_TYPES}"
            )

    theme_value = theme if theme is not None else spec.get("theme")

    if style not in list_recipes():
        raise ValueError(f"style={style!r} invalid; expected one of {list_recipes()}")

    out_path = output_path(output_filename)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    engine = (engine or "pptxgenjs").lower()
    if engine not in ("pptxgenjs", "python-pptx"):
        raise ValueError(f"engine={engine!r} invalid; expected pptxgenjs or python-pptx")

    effective_spec = dict(spec)
    if theme_value is not None:
        effective_spec["theme"] = theme_value
    # WP1: the pptxgenjs engine never received `style` before — wire it in so
    # sharp/soft/rounded/pill actually drive the Node renderer's visuals.
    effective_spec["style"] = style

    # Cross-cutting: PPTX has no runtime font fallback — a font that isn't
    # installed in the render environment silently degrades (esp. CJK). Warn
    # and substitute a known-safe family. Never blocks rendering.
    font_warnings = _validate_and_fix_fonts(effective_spec)

    if engine == "pptxgenjs":
        from . import pptxgenjs_engine

        result = pptxgenjs_engine.build(effective_spec, output_path=out_path)
        # Diagnose layout variety against the engine's ground-truth per-slide
        # layouts (returned in result["layouts"]) — never re-derive them here.
        layout_warnings = _diagnose_layouts(result.get("layouts") or [])
        return {
            "output_filename": output_filename,
            "engine": "pptxgenjs",
            "theme": _theme_summary(theme_value),
            "style": style,
            "slide_count": int(result.get("slides", len(slides))),
            "size_bytes": out_path.stat().st_size,
            "size_kb": int(result.get("size_kb", 0)) or round(out_path.stat().st_size / 1024),
            "layout_warnings": layout_warnings,
            "font_warnings": font_warnings,
        }

    # python-pptx engine (no-Node backup) — the reduced renderer degrades rich
    # layouts to bullets, so variety diagnostics don't apply.
    layout_warnings: list[str] = []
    palette = _resolve_palette_for_python_pptx(theme_value)
    recipe = get_recipe(style)
    fonts = spec.get("fonts") if isinstance(spec.get("fonts"), dict) else None

    prs = _new_presentation()
    for idx, slide_spec in enumerate(slides):
        layout = _blank_layout(prs)  # robust across pptxgenjs + python-pptx masters
        slide = prs.slides.add_slide(layout)
        slide_type = str(slide_spec.get("type") or "content").lower()
        title = slide_spec.get("title")
        content = {k: v for k, v in slide_spec.items() if k not in ("type", "title")}
        slide_types.render(
            slide=slide,
            slide_type=slide_type,
            title=title,
            content=content,
            palette=palette,
            recipe=recipe,
            index=idx + 1,
            fonts=fonts,
        )

    prs.save(str(out_path))
    return {
        "output_filename": output_filename,
        "engine": "python-pptx",
        "theme": _theme_summary(theme_value),
        "style": style,
        "slide_count": len(slides),
        "size_bytes": out_path.stat().st_size,
        "layout_warnings": layout_warnings,
        "font_warnings": font_warnings,
    }


# ── Public: single-slide editing (preserved for editing existing decks) ─


def add_slide(
    *,
    input_filename: str,
    output_filename: str,
    slide_type: str,
    title: str | None = None,
    content: dict[str, Any] | None = None,
    theme: str = "default",
    style: str = "soft",
    index: int | None = None,
    fonts: dict | None = None,
) -> dict[str, Any]:
    """Append a slide to an existing presentation.

    Args:
        input_filename: source .pptx in workdir
        output_filename: destination .pptx
        slide_type: one of ``cover``, ``toc``, ``section``, ``content``, ``summary``
        title: slide title (semantics vary by slide_type)
        content: extra fields per slide_type (subtitle / items / bullets /
            leftBullets+rightBullets / highlights / body / author)
        theme: palette name or alias (independent of any theme baked into the
            source deck — we use absolute colors)
        style: visual style — ``sharp`` / ``soft`` / ``rounded`` / ``pill``
        index: 1-based page badge number; defaults to the new slide's
            position (which is the natural choice for an appended slide)

    Returns:
        ``{output_filename, slide_index, slide_type, title, slide_count, ...}``
    """
    from pptx import Presentation

    if slide_type not in VALID_SLIDE_TYPES:
        raise ValueError(
            f"slide_type={slide_type!r} invalid; expected one of {VALID_SLIDE_TYPES}"
        )

    prs = Presentation(str(input_path(input_filename)))
    palette = get_palette(theme)
    recipe = get_recipe(style)

    layout = _blank_layout(prs)  # robust across pptxgenjs + python-pptx masters
    slide = prs.slides.add_slide(layout)
    new_index = len(prs.slides)  # 1-based count of slides AFTER append
    badge_index = index if index is not None else new_index

    meta = slide_types.render(
        slide=slide,
        slide_type=slide_type,
        title=title,
        content=content,
        palette=palette,
        recipe=recipe,
        index=badge_index,
        fonts=fonts,
    )

    out = output_path(output_filename)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    return {
        "output_filename": output_filename,
        "slide_index": new_index - 1,
        "slide_count": new_index,
        "theme": theme,
        "style": style,
        **meta,
    }


# ── helpers ──


def _diagnose_layouts(layouts: list[dict[str, Any]]) -> list[str]:
    """Spot lazy / repetitive layout choices and surface as warnings.

    Consumes the engine's ground-truth per-slide layouts (the ``layouts`` array
    returned by node_scripts/build_presentation.js — ``{index, type, layout}``
    per slide), so the diagnostic can never drift from what was actually
    rendered. Returns human-readable strings the LLM should react to by
    re-calling build_from_spec with a richer spec. Empty list = clean.
    """
    warnings: list[str] = []
    content_slides = [
        (int(d.get("index", i + 1)), str(d.get("layout") or "single"))
        for i, d in enumerate(layouts)
        if str(d.get("type", "content")).lower() == "content"
    ]
    n_content = len(content_slides)

    # Rule 1: 3+ consecutive plain `single` (bullets-only) slides
    run_start: int | None = None
    run_len = 0
    for idx, layout in content_slides:
        if layout == "single":
            if run_start is None:
                run_start = idx
            run_len += 1
        else:
            if run_len >= 3:
                warnings.append(
                    f"slides {run_start}-{run_start + run_len - 1} all use plain bullets "
                    f"({run_len} in a row) — replace with icon_rows / stat_callout / "
                    f"grid / timeline so the deck has visual variety"
                )
            run_start = None
            run_len = 0
    if run_len >= 3:
        warnings.append(
            f"slides {run_start}-{run_start + run_len - 1} all use plain bullets "
            f"({run_len} in a row) — replace with icon_rows / stat_callout / "
            f"grid / timeline so the deck has visual variety"
        )

    # Rule 2: deck has many content slides but no layout variety
    if n_content >= 4:
        unique_layouts = {layout for _, layout in content_slides}
        if len(unique_layouts) == 1 and "single" in unique_layouts:
            warnings.append(
                f"all {n_content} content slides use plain bullets — pick at least "
                f"3 different layouts (stat_callout for KPIs, icon_rows for "
                f"features, grid for categories, timeline for steps)"
            )
        elif len(unique_layouts) < 3 and n_content >= 6:
            warnings.append(
                f"only {len(unique_layouts)} layout(s) used across {n_content} "
                f"content slides — mix in at least 3 different layouts"
            )

    return warnings


def _theme_summary(theme_value: Any) -> Any:
    """Compact theme description for the response envelope."""
    if theme_value is None:
        return "default"
    if isinstance(theme_value, str):
        return theme_value
    if isinstance(theme_value, dict):
        return "<custom>"
    return str(theme_value)


_SAFE_FONT_CHAIN = (
    "Microsoft YaHei", "Noto Sans CJK SC", "Source Han Sans SC",
    "WenQuanYi Zen Hei", "WenQuanYi Micro Hei", "SimHei",
    "Noto Sans SC", "Arial",
)


def _installed_font_families() -> set[str] | None:
    """Lowercased set of font families known to fontconfig.

    Returns None when fc-list is unavailable (dev box without fontconfig) —
    callers then skip validation rather than emit false warnings.
    """
    import shutil
    import subprocess

    if not shutil.which("fc-list"):
        return None
    try:
        out = subprocess.run(
            ["fc-list", ":", "family"],
            capture_output=True, text=True, timeout=10,
        )
    except Exception:
        return None
    if out.returncode != 0:
        return None
    fams: set[str] = set()
    for line in out.stdout.splitlines():
        for fam in line.split(","):
            fam = fam.strip()
            if fam:
                fams.add(fam.lower())
    return fams or None


def _validate_and_fix_fonts(effective_spec: dict[str, Any]) -> list[str]:
    """Warn + substitute when requested fonts aren't installed (no-op if
    fontconfig is unavailable). Mutates ``effective_spec['fonts']``."""
    installed = _installed_font_families()
    if installed is None:
        return []

    # The Node engine defaults to "Microsoft YaHei" when fonts is unset —
    # validate the effective values so missing CJK on Linux is caught too.
    fonts = effective_spec.get("fonts")
    fonts = dict(fonts) if isinstance(fonts, dict) else {}
    requested = {
        "header": fonts.get("header") or "Microsoft YaHei",
        "body": fonts.get("body") or "Microsoft YaHei",
    }

    safe = next((f for f in _SAFE_FONT_CHAIN if f.lower() in installed), None)
    warnings: list[str] = []
    for role, fam in requested.items():
        if fam.lower() in installed:
            fonts[role] = fam
            continue
        if safe:
            fonts[role] = safe
            warnings.append(
                f"font {fam!r} ({role}) is not installed in the render "
                f"environment — substituted {safe!r} to avoid silent "
                f"degradation (esp. CJK)."
            )
        else:
            warnings.append(
                f"font {fam!r} ({role}) is not installed and no safe CJK "
                f"fallback was found; output may render with substituted glyphs."
            )
    if warnings:
        effective_spec["fonts"] = fonts
    return warnings


def _resolve_palette_for_python_pptx(theme_value: Any):
    """Convert spec.theme (string | dict | None) → Palette TypedDict."""
    if theme_value is None:
        return get_palette(None)
    if isinstance(theme_value, str):
        return get_palette(theme_value)
    if isinstance(theme_value, dict):
        base = dict(get_palette(None))
        for key in ("primary", "secondary", "accent", "light", "bg", "text_on_primary"):
            v = theme_value.get(key)
            if isinstance(v, str) and v.strip():
                base[key] = v.strip().lstrip("#").upper()
        return base  # type: ignore[return-value]
    return get_palette(None)
