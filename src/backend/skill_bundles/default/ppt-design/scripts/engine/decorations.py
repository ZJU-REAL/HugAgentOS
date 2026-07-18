"""Reusable decorative primitives for python-pptx slides.

These helpers wrap python-pptx's ``slides.shapes.add_shape`` and ``add_textbox``
to draw the visual elements that turn a plain text-only slide into a
presentable one — accent bars, page badges, highlight cards, icon circles,
shadows. All renderers in ``slide_types.py`` should call into this module
rather than re-implementing shapes inline.

Coordinate convention: all ``x/y/w/h`` are in INCHES. Internally converted
to EMU via ``pptx.util.Inches``. The 16:9 canvas is 10" × 5.625".

⚠️ Note: ``add_divider_line`` is intentionally NOT exposed. The pptx-skill
guidelines explicitly call out title underlines as a "hallmark of
AI-generated slides" — use whitespace, background-color zoning, or accent
dots instead.
"""
from __future__ import annotations

from typing import Any

from ._colors import hex_to_rgb
from .themes import Palette
from .style_recipes import StyleRecipe
from . import typography


def _solid_fill(shape, hex_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(hex_color)


def _no_line(shape) -> None:
    shape.line.fill.background()


def _solid_line(shape, hex_color: str, weight_pt: float = 1.0) -> None:
    from pptx.util import Pt
    shape.line.color.rgb = hex_to_rgb(hex_color)
    shape.line.width = Pt(weight_pt)


def add_accent_bar(slide, *, x: float, y: float, w: float, h: float, color_hex: str):
    """Solid filled rectangle — used as cover left bar / section accent block."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    _solid_fill(shape, color_hex)
    _no_line(shape)
    return shape


def add_accent_dot(slide, *, x: float, y: float, diameter: float, color_hex: str):
    """Filled circle — primary visual motif. Use to mark covers / sections."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(diameter), Inches(diameter),
    )
    _solid_fill(shape, color_hex)
    _no_line(shape)
    return shape


def add_icon_circle(
    slide,
    *,
    x: float,
    y: float,
    diameter: float,
    fill_hex: str,
    glyph: str,
    glyph_color: str = "FFFFFF",
    glyph_size_pt: int | None = None,
):
    """Filled circle with a single unicode glyph centered inside.

    Used as the icon column of the ``icon_rows`` content layout. We rely on
    unicode geometric symbols (●○◆◇▲△■□★▶) instead of pulling in a heavy
    react-icons + sharp dependency.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(diameter), Inches(diameter),
    )
    _solid_fill(circle, fill_hex)
    _no_line(circle)

    tf = circle.text_frame
    tf.margin_left = tf.margin_right = Inches(0.0)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = str(glyph)
    run.font.size = Pt(glyph_size_pt or max(int(diameter * 36), 12))  # scale to circle
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(glyph_color)
    return circle


def add_page_badge(slide, *, index: int, accent_hex: str, recipe: StyleRecipe, text_color: str = "FFFFFF"):
    """Auto page-number badge at bottom-right (rounded rectangle + white digits).

    Position fixed at (9.08", 5.08", 0.62" × 0.34") — same coords used by
    ``node_scripts/build_presentation.js``.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.08), Inches(5.08), Inches(0.62), Inches(0.34),
    )
    _solid_fill(shape, accent_hex)
    _no_line(shape)
    try:
        target_radius = min(recipe["radius_medium"], 0.34 / 2)
        shape.adjustments[0] = max(0.0, min(0.5, target_radius / 0.34))
    except Exception:
        pass

    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = f"{index:02d}"
    run.font.size = Pt(typography.CAPTION)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(text_color)
    return shape


def add_highlight_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    fill_hex: str,
    border_hex: str | None,
    text: str,
    text_color: str,
    recipe: StyleRecipe,
    font_size: int | None = None,
    bold: bool = True,
    add_shadow: bool = False,
):
    """Rounded rectangle with centered single-line text — bottom highlight pills."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    _solid_fill(shape, fill_hex)
    if border_hex:
        _solid_line(shape, border_hex, weight_pt=0.75)
    else:
        _no_line(shape)
    try:
        shape.adjustments[0] = max(0.0, min(0.5, recipe["radius_medium"] / max(h, 0.001)))
    except Exception:
        pass

    if add_shadow:
        _apply_shadow(shape, blur_pt=8, offset_pt=2, opacity=0.18)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(recipe["padding_min"])
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size or typography.CAPTION)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(text_color)
    return shape


def add_floating_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    fill_hex: str = "FFFFFF",
    border_hex: str | None = None,
    recipe: StyleRecipe | None = None,
    add_shadow: bool = False,
):
    """Empty rounded card used as a body container floating over a tinted bg."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    _solid_fill(shape, fill_hex)
    if border_hex:
        _solid_line(shape, border_hex, weight_pt=0.75)
    else:
        _no_line(shape)
    if recipe is not None:
        try:
            shape.adjustments[0] = max(0.0, min(0.5, recipe["radius_small"] / max(h, 0.001)))
        except Exception:
            pass
    if add_shadow:
        _apply_shadow(shape, blur_pt=10, offset_pt=3, opacity=0.20)
    try:
        shape.text_frame.text = ""
    except Exception:
        pass
    return shape


def fill_slide_background(slide, color_hex: str) -> None:
    """Set the entire slide background to a solid color."""
    from pptx.dml.color import RGBColor

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(color_hex.lstrip("#"))


# ── Shadow (raw OOXML — python-pptx has no high-level API) ─────────────


def _apply_shadow(shape, *, blur_pt: float = 8, offset_pt: float = 2, opacity: float = 0.20, angle_deg: int = 135) -> None:
    """Apply a soft outer shadow to ``shape`` via raw spPr XML manipulation.

    python-pptx has no high-level shadow API. We splice an ``<a:outerShdw>``
    element into the shape's ``<a:effectLst>``. EMU conversions: 1pt = 12700,
    angles are 60000-units (135° × 60000 = 8100000), opacity is 100000-units.
    """
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    sp_pr = shape._element.spPr
    # Drop any existing effectLst before injecting a new one
    for existing in sp_pr.findall("a:effectLst", nsmap):
        sp_pr.remove(existing)

    blur_emu = int(blur_pt * 12700)
    offset_emu = int(offset_pt * 12700)
    angle_units = (angle_deg % 360) * 60000
    opacity_units = int(max(0.0, min(1.0, opacity)) * 100000)

    effect_xml = (
        f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur_emu}" dist="{offset_emu}" '
        f'dir="{angle_units}" algn="tl" rotWithShape="0">'
        f'<a:srgbClr val="000000"><a:alpha val="{opacity_units}"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst>'
    )
    effect = etree.fromstring(effect_xml)
    sp_pr.append(effect)


__all__ = [
    "add_accent_bar",
    "add_accent_dot",
    "add_icon_circle",
    "add_page_badge",
    "add_highlight_card",
    "add_floating_card",
    "fill_slide_background",
]
