"""Slide type renderers for python-pptx — visually aligned with skill guidance.

Layout catalog by slide_type:

    cover   — dark primary bg + small accent dot motif + large title +
              optional italic tagline + meta line. Sandwich opener.

    toc     — light bg + numbered list (NO divider line under title — that
              underline is the AI-generated cliché the pptx skill warns about)

    section — primary bg + LEFT half-bleed accent panel (optional) + accent
              dot motif + big title with charSpacing + optional subtitle.

    content — title + body. Sub-layout chosen by content fields:
        - layout="icon_rows"     items: [{glyph, title, desc}]  (3-5 rows)
        - layout="stat_callout"  stats: [{value, label, tagline?}]  (1-3 cols)
        - layout="grid"          items: [{title, desc}]          (auto 2x2/2x3)
        - layout="timeline"      steps: [{step, title, desc}]    (3-5 horiz)
        - layout="two_col"       leftBullets / rightBullets        (compare)
        - layout="highlights"    bullets + highlights (≤3 pills)
        - layout="single"        bullets                            (default)

    summary — dark primary bg + floating white card with shadow +
              checkmark bullets. Sandwich closer.

The renderers use python-pptx's ``slides.add_slide(layout)`` with absolute
positioning rather than the master's text placeholders — consistent
rendering regardless of which template the deck started from.

Three knobs all renderers share:
- ``palette``  — color scheme (see ``themes.py``)
- ``recipe``   — style recipe controlling corner radius + spacing
- ``index``    — 1-based slide number for the auto page badge
"""
from __future__ import annotations

from typing import Any

from . import decorations as deco
from . import typography as typ
from ._colors import hex_to_rgb, set_char_spacing
from .style_recipes import StyleRecipe, get_recipe
from .themes import Palette


_DEFAULT_HEADER_FONT = "Microsoft YaHei"
_DEFAULT_BODY_FONT = "Microsoft YaHei"


def _set_paragraph_font(
    para,
    *,
    size_pt: int,
    bold: bool = False,
    italic: bool = False,
    color_hex: str | None = None,
    font_face: str | None = None,
    char_spacing_pt: float | None = None,
) -> None:
    from pptx.util import Pt

    para.font.size = Pt(size_pt)
    para.font.bold = bold
    para.font.italic = italic
    if color_hex:
        para.font.color.rgb = hex_to_rgb(color_hex)
    if font_face:
        para.font.name = font_face
    for run in para.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        if color_hex:
            run.font.color.rgb = hex_to_rgb(color_hex)
        if font_face:
            run.font.name = font_face
        if char_spacing_pt is not None:
            set_char_spacing(run, char_spacing_pt)


def _add_textbox(slide, *, left_in: float, top_in: float, width_in: float, height_in: float, text: str = ""):
    from pptx.util import Inches

    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    tf = box.text_frame
    tf.word_wrap = True
    if text:
        tf.text = text
    return box, tf


def _add_title(
    slide,
    *,
    title: str,
    palette: Palette,
    fonts: dict,
    top_in: float = 0.42,
    size_pt: int | None = None,
    color_hex: str | None = None,
) -> None:
    """Standard page title — used by toc/content/summary. NO divider line below."""
    _, tf = _add_textbox(slide, left_in=0.72, top_in=top_in, width_in=8.2, height_in=0.7, text=title)
    _set_paragraph_font(
        tf.paragraphs[0],
        size_pt=size_pt or typ.TITLE,
        bold=True,
        color_hex=color_hex or palette["primary"],
        font_face=fonts.get("header", _DEFAULT_HEADER_FONT),
    )


def _add_bullets(
    slide,
    *,
    bullets,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    palette: Palette,
    fonts: dict,
    size_pt: int | None = None,
    color_hex: str | None = None,
) -> int:
    """Add a textbox of bullets; returns count of items rendered."""
    items = [str(b) for b in (bullets or []) if b]
    if not items:
        return 0
    _, tf = _add_textbox(slide, left_in=left_in, top_in=top_in, width_in=width_in, height_in=height_in)
    tf.text = ""
    body_color = color_hex or palette["secondary"]
    for i, b in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = "• " + b
        _set_paragraph_font(
            para,
            size_pt=size_pt or typ.BODY + 4,
            color_hex=body_color,
            font_face=fonts.get("body", _DEFAULT_BODY_FONT),
        )
    return len(items)


def _resolve_fonts(spec_fonts: dict | None) -> dict:
    """Normalize the optional ``fonts`` dict from the spec."""
    if not isinstance(spec_fonts, dict):
        return {"header": _DEFAULT_HEADER_FONT, "body": _DEFAULT_BODY_FONT}
    return {
        "header": spec_fonts.get("header") or _DEFAULT_HEADER_FONT,
        "body": spec_fonts.get("body") or _DEFAULT_BODY_FONT,
    }


# ── Public dispatcher ──────────────────────────────────────────────────


def render(
    *,
    slide,
    slide_type: str,
    title: str | None,
    content: dict[str, Any] | None,
    palette: Palette,
    recipe: StyleRecipe | None = None,
    index: int = 1,
    fonts: dict | None = None,
) -> dict[str, Any]:
    """Dispatch to the right renderer based on ``slide_type``.

    ``recipe`` defaults to ``soft`` if not supplied.
    ``index`` is the 1-based slide number for the auto page badge.
    ``fonts`` is the spec-level fonts dict ``{header, body}``; when None,
    Microsoft YaHei is used for both (Chinese-friendly default).
    """
    content = content or {}
    recipe = recipe or get_recipe("soft")
    fonts_resolved = _resolve_fonts(fonts)

    if slide_type == "cover":
        return _render_cover(slide, title or "", content, palette, recipe, fonts_resolved)
    if slide_type == "toc":
        return _render_toc(slide, title or "目录", content, palette, recipe, index, fonts_resolved)
    if slide_type == "section":
        return _render_section(slide, title or "", content, palette, recipe, index, fonts_resolved)
    if slide_type == "summary":
        return _render_summary(slide, title or "总结", content, palette, recipe, index, fonts_resolved)
    if slide_type == "content":
        return _render_content(slide, title or "", content, palette, recipe, index, fonts_resolved)
    if slide_type in ("closing", "thanks"):
        return _render_closing(slide, title or "谢谢观看", content, palette, recipe, fonts_resolved)
    raise ValueError(
        f"unknown slide_type {slide_type!r}; expected one of: "
        "cover, toc, section, content, summary, closing"
    )


# ── Renderers ──────────────────────────────────────────────────────────


def _cover_colors(palette: Palette, style: str) -> dict[str, str]:
    if style == "light":
        return {
            "bg": palette["bg"],
            "title": palette["primary"],
            "sub": palette["secondary"],
            "meta": palette["secondary"],
            "dot": palette["accent"],
            "bar": palette["accent"],
        }
    return {
        "bg": palette["primary"],
        "title": palette["text_on_primary"],
        "sub": palette["light"],
        "meta": palette["light"],
        "dot": palette["accent"],
        "bar": palette["accent"],
    }


def _render_cover(slide, title: str, content: dict, palette: Palette, recipe: StyleRecipe, fonts: dict) -> dict:
    """Cover: light bg + accent dot + large title + optional italic tagline +
    meta line. Light is the default (政府汇报风格); ``content.cover_style="dark"``
    flips to a primary-bg cover for high-impact / brand decks.
    """
    style = str(content.get("cover_style") or "light").lower()
    c = _cover_colors(palette, style)

    deco.fill_slide_background(slide, c["bg"])

    # Visual motif: small accent dot at top-left (carried across cover/section)
    deco.add_accent_dot(slide, x=0.62, y=0.62, diameter=0.22, color_hex=c["dot"])

    # Subtle accent bar to the LEFT of the title block (≠ the rejected "title underline")
    deco.add_accent_bar(slide, x=0.6, y=1.85, w=0.10, h=2.4, color_hex=c["bar"])

    # Big title
    _, tf = _add_textbox(
        slide,
        left_in=0.85, top_in=1.6, width_in=8.2, height_in=1.5,
        text=title,
    )
    _set_paragraph_font(
        tf.paragraphs[0],
        size_pt=typ.LARGE_TITLE,
        bold=True,
        color_hex=c["title"],
        font_face=fonts["header"],
    )

    subtitle = content.get("subtitle") or ""
    if subtitle:
        _, tf2 = _add_textbox(
            slide,
            left_in=0.88, top_in=3.05, width_in=7.6, height_in=0.7,
            text=str(subtitle),
        )
        _set_paragraph_font(
            tf2.paragraphs[0],
            size_pt=typ.SUBTITLE,
            color_hex=c["sub"],
            font_face=fonts["body"],
        )

    tagline = content.get("tagline")
    if tagline:
        _, tf_tag = _add_textbox(
            slide,
            left_in=0.88, top_in=3.85, width_in=7.6, height_in=0.5,
            text=str(tagline),
        )
        _set_paragraph_font(
            tf_tag.paragraphs[0],
            size_pt=typ.SUBTITLE - 2,
            italic=True,
            color_hex=c["sub"],
            font_face=fonts["body"],
        )

    body_or_author = content.get("body") or content.get("author") or ""
    if body_or_author:
        _, tf3 = _add_textbox(
            slide,
            left_in=0.88, top_in=4.75, width_in=4.5, height_in=0.32,
            text=str(body_or_author),
        )
        _set_paragraph_font(
            tf3.paragraphs[0],
            size_pt=typ.CAPTION,
            color_hex=c["meta"],
            font_face=fonts["body"],
        )

    return {"slide_type": "cover", "title": title, "subtitle": subtitle, "cover_style": style}


def _render_toc(slide, title: str, content: dict, palette: Palette, recipe: StyleRecipe, index: int, fonts: dict) -> dict:
    """TOC: light bg + title (no underline) + accent dot + numbered list + page badge."""
    deco.fill_slide_background(slide, "FFFFFF")
    # Dot top-right so it doesn't collide with the left-anchored title
    deco.add_accent_dot(slide, x=9.10, y=0.55, diameter=0.18, color_hex=palette["accent"])
    _add_title(slide, title=title, palette=palette, fonts=fonts, top_in=0.42)

    items = content.get("items") or content.get("bullets") or []
    items = [str(it) for it in items if it]

    _, tf = _add_textbox(slide, left_in=0.95, top_in=1.45, width_in=8.0, height_in=3.4)
    tf.text = ""
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"{i + 1:02d}    {it}"
        _set_paragraph_font(
            para,
            size_pt=typ.SUBTITLE,
            color_hex=palette["secondary"],
            font_face=fonts["body"],
        )
        para.space_after = _pt_to_emu(10)

    deco.add_page_badge(slide, index=index, accent_hex=palette["accent"], recipe=recipe)
    return {"slide_type": "toc", "title": title, "item_count": len(items)}


def _render_section(slide, title: str, content: dict, palette: Palette, recipe: StyleRecipe, index: int, fonts: dict) -> dict:
    """Section: primary bg + LEFT half-bleed accent panel (optional) + accent
    dot motif + big title with character spacing + page badge.
    """
    deco.fill_slide_background(slide, palette["primary"])

    # Optional left half-bleed panel — adds a darker block to balance text
    half_bleed = bool(content.get("half_bleed", True))
    if half_bleed:
        deco.add_accent_bar(slide, x=0.0, y=0.0, w=3.2, h=5.625, color_hex=palette["secondary"])

    # Accent dot motif (top-right when half-bleed left, top-left otherwise)
    if half_bleed:
        deco.add_accent_dot(slide, x=9.10, y=0.62, diameter=0.22, color_hex=palette["accent"])
    else:
        deco.add_accent_dot(slide, x=0.78, y=0.55, diameter=0.22, color_hex=palette["accent"])

    title_left = 3.55 if half_bleed else 0.78
    title_width = 6.0 if half_bleed else 8.4

    # Big section title with letter-spacing for premium feel
    _, tf = _add_textbox(
        slide,
        left_in=title_left, top_in=2.0, width_in=title_width, height_in=1.4,
        text=title,
    )
    _set_paragraph_font(
        tf.paragraphs[0],
        size_pt=typ.LARGE_TITLE - 4,
        bold=True,
        color_hex=palette["text_on_primary"],
        font_face=fonts["header"],
        char_spacing_pt=2.0,
    )

    subtitle = content.get("subtitle") or ""
    if subtitle:
        _, tf2 = _add_textbox(
            slide,
            left_in=title_left, top_in=3.45, width_in=title_width, height_in=0.6,
            text=str(subtitle),
        )
        _set_paragraph_font(
            tf2.paragraphs[0],
            size_pt=typ.BODY + 1,
            color_hex=palette["light"],
            font_face=fonts["body"],
        )

    deco.add_page_badge(slide, index=index, accent_hex=palette["light"], recipe=recipe, text_color=palette["primary"])
    return {"slide_type": "section", "title": title, "half_bleed": half_bleed}


def _detect_layout(content: dict) -> str:
    """Pick a content sub-layout based on explicit field or content shape."""
    explicit = content.get("layout")
    if explicit:
        return str(explicit).lower()
    if content.get("items") and (
        all(isinstance(it, dict) and ("glyph" in it or "icon" in it) for it in content["items"])
    ):
        return "icon_rows"
    if content.get("stats"):
        return "stat_callout"
    if content.get("steps"):
        return "timeline"
    if content.get("items") and all(isinstance(it, dict) for it in content["items"]):
        return "grid"
    if content.get("leftBullets") or content.get("rightBullets"):
        return "two_col"
    if content.get("highlights"):
        return "highlights"
    return "single"


def _render_content(slide, title: str, content: dict, palette: Palette, recipe: StyleRecipe, index: int, fonts: dict) -> dict:
    """Content slide — dispatches to one of 7 sub-layouts."""
    deco.fill_slide_background(slide, palette["bg"])
    _add_title(slide, title=title, palette=palette, fonts=fonts)

    layout = _detect_layout(content)
    item_count = 0

    if layout == "icon_rows":
        item_count = _content_icon_rows(slide, content, palette, recipe, fonts)
    elif layout == "stat_callout":
        item_count = _content_stat_callout(slide, content, palette, recipe, fonts)
    elif layout == "grid":
        item_count = _content_grid(slide, content, palette, recipe, fonts)
    elif layout == "timeline":
        item_count = _content_timeline(slide, content, palette, recipe, fonts)
    elif layout == "two_col":
        item_count = _content_two_col(slide, content, palette, recipe, fonts)
    elif layout == "highlights":
        item_count = _content_highlights(slide, content, palette, recipe, fonts)
    else:
        item_count = _content_single(slide, content, palette, recipe, fonts)

    body = content.get("body")
    if body and layout not in ("highlights", "stat_callout"):
        _, tfb = _add_textbox(slide, left_in=0.90, top_in=4.85, width_in=8.15, height_in=0.4, text=str(body))
        _set_paragraph_font(tfb.paragraphs[0], size_pt=typ.CAPTION + 1, color_hex=palette["secondary"], font_face=fonts["body"])

    deco.add_page_badge(slide, index=index, accent_hex=palette["accent"], recipe=recipe)
    return {
        "slide_type": "content",
        "title": title,
        "layout": layout,
        "item_count": item_count,
    }


# ── Content sub-layouts ───────────────────────────────────────────────


def _content_single(slide, content, palette, recipe, fonts) -> int:
    return _add_bullets(
        slide, bullets=content.get("bullets"),
        left_in=0.88, top_in=1.45, width_in=8.25, height_in=3.2,
        palette=palette, fonts=fonts,
    )


def _content_two_col(slide, content, palette, recipe, fonts) -> int:
    """Two-column compare. Subtle vertical accent bar between columns."""
    if content.get("leftTitle"):
        _, tfl = _add_textbox(slide, left_in=0.78, top_in=1.40, width_in=3.9, height_in=0.45, text=str(content["leftTitle"]))
        _set_paragraph_font(tfl.paragraphs[0], size_pt=typ.SUBTITLE - 2, bold=True, color_hex=palette["primary"], font_face=fonts["header"])
    if content.get("rightTitle"):
        _, tfr = _add_textbox(slide, left_in=5.30, top_in=1.40, width_in=3.9, height_in=0.45, text=str(content["rightTitle"]))
        _set_paragraph_font(tfr.paragraphs[0], size_pt=typ.SUBTITLE - 2, bold=True, color_hex=palette["primary"], font_face=fonts["header"])

    # Faint vertical separator between columns (not under the title — that's the banned underline)
    deco.add_accent_bar(slide, x=4.95, y=1.55, w=0.04, h=2.9, color_hex=palette["light"])

    left = _add_bullets(slide, bullets=content.get("leftBullets"),
                        left_in=0.80, top_in=1.95, width_in=4.0, height_in=2.6,
                        palette=palette, fonts=fonts)
    right = _add_bullets(slide, bullets=content.get("rightBullets"),
                         left_in=5.30, top_in=1.95, width_in=3.9, height_in=2.6,
                         palette=palette, fonts=fonts)
    return left + right


def _content_highlights(slide, content, palette, recipe, fonts) -> int:
    """Bullets above + 3 highlight cards across bottom."""
    bullet_count = _add_bullets(
        slide, bullets=content.get("bullets"),
        left_in=0.88, top_in=1.45, width_in=8.25, height_in=2.5,
        palette=palette, fonts=fonts,
    )
    hi = content.get("highlights") or []
    hi = [str(h) for h in hi[:3] if h]
    for i, item in enumerate(hi):
        x = 0.88 + i * 2.95
        deco.add_highlight_card(
            slide, x=x, y=4.20, w=2.55, h=0.62,
            fill_hex=palette["light"] if i % 2 else palette["bg"],
            border_hex=palette["light"],
            text=str(item),
            text_color=palette["primary"],
            recipe=recipe,
            add_shadow=True,
        )
    return bullet_count + len(hi)


def _content_icon_rows(slide, content, palette, recipe, fonts) -> int:
    """Icon column + bold header + description row, 3-5 rows.

    Spec field: ``items: [{glyph, title, desc}]`` where ``glyph`` is ONE
    unicode character (e.g. "●", "◆", "▶", "★"). ``icon`` is accepted as a
    synonym for ``glyph`` for LLM-friendliness.
    """
    items = content.get("items") or []
    items = [it for it in items if isinstance(it, dict)][:5]
    if not items:
        return 0

    top = 1.50
    available = 3.4
    row_h = min(0.85, available / len(items))
    gap = 0.10
    icon_diam = min(row_h - 0.05, 0.62)

    for i, it in enumerate(items):
        y = top + i * (row_h + gap)
        glyph = str(it.get("glyph") or it.get("icon") or "●")[:2]
        deco.add_icon_circle(
            slide,
            x=0.88, y=y, diameter=icon_diam,
            fill_hex=palette["accent"],
            glyph=glyph,
            glyph_color=palette["text_on_primary"],
            glyph_size_pt=int(icon_diam * 28),
        )
        title_text = str(it.get("title") or "")
        desc_text = str(it.get("desc") or it.get("description") or "")

        text_left = 0.88 + icon_diam + 0.20
        text_w = 9.0 - text_left - 0.4

        _, tft = _add_textbox(slide, left_in=text_left, top_in=y - 0.02, width_in=text_w, height_in=0.42, text=title_text)
        _set_paragraph_font(tft.paragraphs[0], size_pt=typ.SUBTITLE - 4, bold=True, color_hex=palette["primary"], font_face=fonts["header"])

        if desc_text:
            _, tfd = _add_textbox(slide, left_in=text_left, top_in=y + 0.36, width_in=text_w, height_in=row_h - 0.36, text=desc_text)
            _set_paragraph_font(tfd.paragraphs[0], size_pt=typ.BODY, color_hex=palette["secondary"], font_face=fonts["body"])
    return len(items)


def _content_stat_callout(slide, content, palette, recipe, fonts) -> int:
    """1-3 large stat numbers across the slide.

    Spec field: ``stats: [{value, label, tagline?}]``. Numbers render at
    DATA_CALLOUT (72pt). When 3 stats are given, columns split evenly.
    """
    stats = content.get("stats") or []
    stats = [s for s in stats if isinstance(s, dict)][:3]
    if not stats:
        return 0

    n = len(stats)
    col_w = 8.4 / n
    base_x = 0.80
    top = 1.65

    for i, s in enumerate(stats):
        x = base_x + i * col_w
        value = str(s.get("value") or "")
        label = str(s.get("label") or "")
        tagline = s.get("tagline")

        # Big number
        _, tfv = _add_textbox(slide, left_in=x, top_in=top, width_in=col_w - 0.2, height_in=1.4, text=value)
        _set_paragraph_font(
            tfv.paragraphs[0],
            size_pt=typ.DATA_CALLOUT,
            bold=True,
            color_hex=palette["primary"],
            font_face=fonts["header"],
        )

        # Label
        if label:
            _, tfl = _add_textbox(slide, left_in=x, top_in=top + 1.45, width_in=col_w - 0.2, height_in=0.4, text=label)
            _set_paragraph_font(tfl.paragraphs[0], size_pt=typ.BODY + 1, color_hex=palette["secondary"], font_face=fonts["body"])

        # Italic tagline
        if tagline:
            _, tft = _add_textbox(slide, left_in=x, top_in=top + 1.95, width_in=col_w - 0.2, height_in=0.5, text=str(tagline))
            _set_paragraph_font(
                tft.paragraphs[0],
                size_pt=typ.BODY,
                italic=True,
                color_hex=palette["accent"],
                font_face=fonts["body"],
            )

    return n


def _content_grid(slide, content, palette, recipe, fonts) -> int:
    """Auto 2x2 / 2x3 / 3x2 grid of {title, desc} cards.

    Spec field: ``items: [{title, desc}]`` (4-6 entries). Card layout
    derived from item count: 4→2x2, 6→2x3 / 3x2 (wider preferred).
    """
    items = content.get("items") or []
    # Cap raised 6→9: ranked data lists (e.g. TOP8) must not lose rows in
    # the python-pptx fallback (mirrors the Node engine's data handling).
    items = [it for it in items if isinstance(it, dict)][:9]
    n = len(items)
    if not n:
        return 0

    # Layout: 1-2 single row; 3 → 1x3; 4 → 2x2; 5-6 → 2x3; 7-9 → 3x3
    if n <= 2:
        rows, cols = 1, n
    elif n == 3:
        rows, cols = 1, 3
    elif n == 4:
        rows, cols = 2, 2
    elif n <= 6:
        rows, cols = 2, 3
    else:
        rows, cols = 3, 3

    # Available area: y 1.45 → 4.65, x 0.78 → 9.22
    avail_w = 9.22 - 0.78
    avail_h = 4.65 - 1.45
    gap = 0.20
    card_w = (avail_w - gap * (cols - 1)) / cols
    card_h = (avail_h - gap * (rows - 1)) / rows

    for idx, it in enumerate(items):
        r = idx // cols
        c = idx % cols
        x = 0.78 + c * (card_w + gap)
        y = 1.45 + r * (card_h + gap)
        # Card background
        deco.add_floating_card(
            slide, x=x, y=y, w=card_w, h=card_h,
            fill_hex="FFFFFF",
            border_hex=palette["light"],
            recipe=recipe,
            add_shadow=False,
        )
        # Title
        title_text = str(it.get("title") or "")
        # Fall back to subtitle/value so a model that packs the figure into
        # `value` (no `desc`) doesn't render an empty card.
        desc_text = str(
            it.get("desc") or it.get("description") or it.get("subtitle")
            or (it.get("value") if it.get("value") not in (None, "") else "")
            or ""
        )
        pad = 0.18
        _, tft = _add_textbox(slide, left_in=x + pad, top_in=y + pad, width_in=card_w - 2 * pad, height_in=0.45, text=title_text)
        _set_paragraph_font(tft.paragraphs[0], size_pt=typ.SUBTITLE - 4, bold=True, color_hex=palette["primary"], font_face=fonts["header"])
        if desc_text:
            _, tfd = _add_textbox(
                slide,
                left_in=x + pad, top_in=y + pad + 0.45,
                width_in=card_w - 2 * pad, height_in=card_h - pad - 0.5,
                text=desc_text,
            )
            _set_paragraph_font(tfd.paragraphs[0], size_pt=typ.BODY, color_hex=palette["secondary"], font_face=fonts["body"])

    return n


def _content_timeline(slide, content, palette, recipe, fonts) -> int:
    """Horizontal numbered timeline. Spec: ``steps: [{step, title, desc}]``.

    Renders a horizontal connector line + numbered circles + step titles + descs.
    """
    steps = content.get("steps") or []
    steps = [s for s in steps if isinstance(s, dict)][:5]
    n = len(steps)
    if not n:
        return 0

    # Horizontal layout: y center for connector at 2.55"
    y_center = 2.55
    left = 1.0
    right = 9.0
    avail = right - left
    spacing = avail / max(n - 1, 1) if n > 1 else 0
    diameter = 0.62

    # Connector line behind the circles
    deco.add_accent_bar(
        slide,
        x=left + diameter / 2,
        y=y_center + diameter / 2 - 0.025,
        w=spacing * (n - 1),
        h=0.05,
        color_hex=palette["light"],
    )

    for i, s in enumerate(steps):
        cx = left + i * spacing if n > 1 else left + avail / 2
        step_label = str(s.get("step") or (i + 1))
        title_text = str(s.get("title") or "")
        desc_text = str(s.get("desc") or s.get("description") or "")

        # Numbered circle
        deco.add_icon_circle(
            slide,
            x=cx, y=y_center, diameter=diameter,
            fill_hex=palette["accent"],
            glyph=step_label[:2],
            glyph_color=palette["text_on_primary"],
            glyph_size_pt=int(diameter * 28),
        )

        col_w = max(min(spacing * 0.95, 2.4), 1.2) if n > 1 else 4.5
        col_x = cx + diameter / 2 - col_w / 2

        # Step title above the circle
        _, tft = _add_textbox(slide, left_in=col_x, top_in=y_center - 0.85, width_in=col_w, height_in=0.4, text=title_text)
        para = tft.paragraphs[0]
        _set_paragraph_font(para, size_pt=typ.SUBTITLE - 4, bold=True, color_hex=palette["primary"], font_face=fonts["header"])
        from pptx.enum.text import PP_ALIGN
        para.alignment = PP_ALIGN.CENTER

        # Step description below the circle
        if desc_text:
            _, tfd = _add_textbox(slide, left_in=col_x, top_in=y_center + diameter + 0.10, width_in=col_w, height_in=1.5, text=desc_text)
            para = tfd.paragraphs[0]
            _set_paragraph_font(para, size_pt=typ.CAPTION + 1, color_hex=palette["secondary"], font_face=fonts["body"])
            para.alignment = PP_ALIGN.CENTER

    return n


def _render_closing(slide, title: str, content: dict, palette: Palette, recipe: StyleRecipe, fonts: dict) -> dict:
    """Closing / thank-you page (python-pptx fallback). Brand-filled by default;
    centered title + subtitle + a row of contact lines. Mirrors the Node engine's
    ``renderClosing`` at a simpler fidelity.
    """
    from pptx.enum.text import PP_ALIGN

    style = str(content.get("closing_style") or "dark").lower()
    dark = style != "light"
    bg = palette["primary"] if dark else palette["bg"]
    title_color = palette["text_on_primary"] if dark else palette["primary"]
    sub_color = palette["light"] if dark else palette["secondary"]

    deco.fill_slide_background(slide, bg)
    deco.add_accent_dot(slide, x=4.89, y=1.30, diameter=0.22, color_hex=palette["accent"])

    _, tf = _add_textbox(slide, left_in=0.5, top_in=1.85, width_in=9.0, height_in=1.3, text=title)
    para = tf.paragraphs[0]
    _set_paragraph_font(para, size_pt=typ.LARGE_TITLE, bold=True, color_hex=title_color, font_face=fonts["header"])
    para.alignment = PP_ALIGN.CENTER

    subtitle = content.get("subtitle") or ""
    if subtitle:
        _, tf2 = _add_textbox(slide, left_in=1.0, top_in=3.2, width_in=8.0, height_in=0.5, text=str(subtitle))
        p2 = tf2.paragraphs[0]
        _set_paragraph_font(p2, size_pt=typ.SUBTITLE, color_hex=sub_color, font_face=fonts["body"])
        p2.alignment = PP_ALIGN.CENTER

    contacts = [c for c in (content.get("contact") or []) if c]
    if contacts:
        parts = []
        for c in contacts:
            if isinstance(c, dict):
                lbl = str(c.get("label") or "")
                val = str(c.get("value") or c.get("label") or "")
                parts.append(f"{lbl}：{val}" if lbl and val and lbl != val else val)
            else:
                parts.append(str(c))
        _, tf3 = _add_textbox(slide, left_in=0.8, top_in=4.15, width_in=8.4, height_in=0.5,
                              text="      ".join(parts))
        p3 = tf3.paragraphs[0]
        _set_paragraph_font(p3, size_pt=typ.CAPTION + 1, color_hex=sub_color, font_face=fonts["body"])
        p3.alignment = PP_ALIGN.CENTER

    return {"slide_type": "closing", "title": title, "closing_style": style}


def _summary_colors(palette: Palette, style: str) -> dict[str, Any]:
    if style == "light":
        return {
            "bg": palette["bg"],
            "title": palette["primary"],
            "bullet": palette["secondary"],
            "body": palette["secondary"],
            "badge_text": "FFFFFF",
            "badge_accent": palette["accent"],
            "card_border": palette["light"],
            "dot_x": 9.10,  # top-right; title is on the left, avoid overlap
        }
    return {
        "bg": palette["primary"],
        "title": palette["text_on_primary"],
        "bullet": palette["primary"],
        "body": palette["secondary"],
        "badge_text": palette["primary"],
        "badge_accent": palette["light"],
        "card_border": None,
        "dot_x": 9.10,
    }


def _render_summary(slide, title: str, content: dict, palette: Palette, recipe: StyleRecipe, index: int, fonts: dict) -> dict:
    """Summary: light bg + floating white card with light border + shadow.
    Light is the default (政府汇报风格); ``content.summary_style="dark"``
    flips to a primary-bg closer to mirror a dark cover.
    """
    style = str(content.get("summary_style") or "light").lower()
    c = _summary_colors(palette, style)

    deco.fill_slide_background(slide, c["bg"])
    _add_title(slide, title=title, palette=palette, fonts=fonts, top_in=0.52, color_hex=c["title"])
    deco.add_accent_dot(slide, x=c["dot_x"], y=0.55, diameter=0.18, color_hex=palette["accent"])

    deco.add_floating_card(
        slide,
        x=0.76, y=1.45, w=8.48, h=3.30,
        fill_hex="FFFFFF",
        border_hex=c["card_border"],
        recipe=recipe,
        add_shadow=True,
    )

    bullets = [str(b) for b in (content.get("bullets") or content.get("highlights") or []) if b]
    if bullets:
        _, tf = _add_textbox(slide, left_in=1.05, top_in=1.75, width_in=7.9, height_in=2.30)
        tf.text = ""
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = "✓  " + b
            _set_paragraph_font(para, size_pt=typ.SUBTITLE, color_hex=c["bullet"], font_face=fonts["body"])
            para.space_after = _pt_to_emu(8)

    body = content.get("body")
    if body:
        _, tfb = _add_textbox(slide, left_in=1.06, top_in=4.20, width_in=7.8, height_in=0.40, text=str(body))
        _set_paragraph_font(tfb.paragraphs[0], size_pt=typ.CAPTION, italic=True, color_hex=c["body"], font_face=fonts["body"])

    deco.add_page_badge(slide, index=index, accent_hex=c["badge_accent"], recipe=recipe, text_color=c["badge_text"])
    return {"slide_type": "summary", "title": title, "summary_style": style, "bullet_count": len(bullets)}


# ── tiny utility ──


def _pt_to_emu(pt: float) -> int:
    """Points → EMU (1 pt = 12700 EMU)."""
    return int(pt * 12700)
