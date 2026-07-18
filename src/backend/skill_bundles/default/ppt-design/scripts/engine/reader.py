"""Read operations on .pptx — slide count, per-slide text content."""
from __future__ import annotations

from typing import Any

from ._shims import input_path


def get_slide_count(*, input_filename: str) -> dict[str, Any]:
    """Return the slide count of a presentation."""
    from pptx import Presentation

    prs = Presentation(str(input_path(input_filename)))
    return {"slide_count": len(prs.slides)}


def get_slide_content(*, input_filename: str, slide_index: int) -> dict[str, Any]:
    """Extract text content from a single slide.

    Args:
        input_filename: source .pptx in cwd
        slide_index: 0-based slide index

    Returns:
        ``{"slide_index", "shape_count", "text_blocks": [{"shape", "text"}, ...],
            "joined_text"}``
    """
    from pptx import Presentation

    prs = Presentation(str(input_path(input_filename)))
    if not (0 <= slide_index < len(prs.slides)):
        raise ValueError(
            f"slide_index {slide_index} out of range "
            f"(presentation has {len(prs.slides)} slides)"
        )
    slide = prs.slides[slide_index]

    text_blocks: list[dict[str, Any]] = []
    joined: list[str] = []
    for i, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        chunks: list[str] = []
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs)
            if line.strip():
                chunks.append(line)
        if chunks:
            txt = "\n".join(chunks)
            text_blocks.append({"shape": i, "text": txt})
            joined.append(txt)

    return {
        "slide_index": slide_index,
        "shape_count": len(slide.shapes),
        "text_blocks": text_blocks,
        "joined_text": "\n\n".join(joined),
    }
