#!/usr/bin/env python3
"""ppt-design skill CLI — single entry point that exposes the pptx engine as
subcommands.

The agent invokes it as ``ppt-cli <subcommand> …`` via Bash — the
``ppt-cli`` shim (``scripts/ppt-cli``) is installed in ``/usr/local/bin/``
inside the sandbox and mcp container images and routes to this file. Local
dev can run ``python scripts/ppt.py …`` directly; behaviour is identical.

Subcommands::

    ppt-cli build --spec spec.json --output deck.pptx [--theme navy_gold --pack dark_gold] [--style soft]
    ppt-cli build-js --script deck.js --output deck.pptx [--timeout 90]
    ppt-cli info deck.pptx
    ppt-cli slide-count deck.pptx
    ppt-cli extract deck.pptx [--slide N]
    ppt-cli check-placeholders deck.pptx [--patterns p1 p2 ...]
    ppt-cli thumbnails deck.pptx --output-dir thumbs/ [--dpi 120] [--quality 85]
    ppt-cli add-slide deck.pptx --type section --output out.pptx [--title …] [--content …] [--theme …] [--style …]
    ppt-cli set-title deck.pptx --slide N --title TXT --output out.pptx
    ppt-cli add-text deck.pptx --slide N --text TXT --output out.pptx [--font-size 14] [--color HEX] [--bold] [--position L,T[,W,H]]
    ppt-cli insert-image deck.pptx --slide N --image IMG.png --output out.pptx [--position L,T[,W,H]]
    ppt-cli delete-slide deck.pptx --slide N --output out.pptx
    ppt-cli to-pdf deck.pptx --output out.pdf
    ppt-cli list-themes
    ppt-cli list-styles
    ppt-cli list-slide-types

Output: on success the CLI prints a single JSON object to stdout describing
the result (output path, meta fields). On failure it writes a JSON error
object to stderr and exits non-zero.

All paths in arguments can be relative or absolute; the CLI normalises them
and copies files in/out of a per-call temp workdir so the engine modules
keep operating on bare filenames.
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, NoReturn

# Make ``engine`` importable when this file is run as a script.
_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))

# Cheap imports only — the build / edit / render submodules pull in python-pptx
# (and transitively lxml + PIL), so we lazy-import them inside each cmd_* that
# needs them. ``list-themes`` / ``list-styles`` / ``list-slide-types`` etc.
# stay snappy as a result.
from engine.style_recipes import list_recipes  # noqa: E402
from engine.themes import list_aliases, list_palettes  # noqa: E402

# These are constants — needed for argparse ``choices=`` at parser-build time,
# and used by ``cmd_list_slide_types``. Defined as module-level tuples here
# rather than re-imported from ``engine.builder`` so the parser stays cheap.
_VALID_ENGINES = ("pptxgenjs", "python-pptx")
_VALID_STYLES = tuple(list_recipes())
_VALID_SLIDE_TYPES = ("cover", "toc", "section", "content", "summary")


# ── helpers ─────────────────────────────────────────────────────────


def _emit(payload: dict[str, Any]) -> None:
    """Print result JSON to stdout — what the agent reads back."""
    sys.stdout.write(json.dumps(payload, ensure_ascii=False))
    sys.stdout.write("\n")


def _die(error_type: str, message: str, *, field: str | None = None, **extra: Any) -> NoReturn:
    """Write an error JSON to stderr and exit with status 2.

    Shape is intentionally aligned with ``mcp_servers/_office_shared/handle.py
    ::error_response`` so agents that already know the MCP error contract
    parse skill-CLI errors with the same code path.
    """
    err: dict[str, Any] = {"ok": False, "error": {"type": error_type, "message": message}}
    if field is not None:
        err["error"]["field"] = field
    if extra:
        err["error"].update(extra)
    sys.stderr.write(json.dumps(err, ensure_ascii=False))
    sys.stderr.write("\n")
    sys.exit(2)


def _ensure_input(path_arg: str) -> Path:
    p = Path(path_arg).expanduser().resolve()
    if not p.is_file():
        _die("FileNotFound", f"input file not found: {p}", field="input")
    return p


def _ensure_pptx_output(path_arg: str) -> Path:
    p = Path(path_arg).expanduser().resolve()
    if p.suffix.lower() != ".pptx":
        _die("ValueError", f"output must end with .pptx, got {p.name!r}", field="output")
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _ensure_pdf_output(path_arg: str) -> Path:
    p = Path(path_arg).expanduser().resolve()
    if p.suffix.lower() != ".pdf":
        _die("ValueError", f"output must end with .pdf, got {p.name!r}", field="output")
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _parse_position(raw: str | None) -> dict[str, float] | None:
    """Accept ``L,T`` or ``L,T,W,H`` (inches) → dict; pass None through."""
    if not raw:
        return None
    parts = [s.strip() for s in raw.split(",") if s.strip()]
    if len(parts) not in (2, 4):
        _die("ValueError", f"--position takes 2 or 4 comma-separated numbers, got {raw!r}", field="position")
    try:
        nums = [float(p) for p in parts]
    except ValueError:
        _die("ValueError", f"--position values must be numbers, got {raw!r}", field="position")
    keys = ("left", "top", "width", "height")
    return {k: v for k, v in zip(keys, nums)}


def _load_spec(spec_arg: str) -> dict[str, Any]:
    """Spec can be either a path to a .json file or an inline JSON string."""
    path = Path(spec_arg).expanduser()
    if path.is_file():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            _die("JSONDecodeError", f"--spec file is not valid JSON: {exc}", field="spec", path=str(path))
    # treat as inline JSON
    try:
        return json.loads(spec_arg)
    except json.JSONDecodeError as exc:
        _die(
            "ValueError",
            f"--spec is neither an existing file path nor valid JSON: {exc}",
            field="spec",
        )


# ── subcommand implementations ──────────────────────────────────────


def cmd_build(args: argparse.Namespace) -> None:
    from engine import _shims, builder

    spec = _load_spec(args.spec)
    if not isinstance(spec, dict):
        _die("ValueError", "spec must be a JSON object (dict)", field="spec")
    out_path = _ensure_pptx_output(args.output)

    # ``--engine`` / ``--style`` are enforced by argparse ``choices=``; only
    # ``--theme`` needs a runtime check because the catalog is dynamic.
    valid_themes = set(list_palettes()) | set(list_aliases().keys())
    if args.theme and args.theme not in valid_themes:
        _die(
            "ValueError",
            f"unknown theme {args.theme!r}; run `ppt-cli list-themes` for the catalog",
            field="theme",
        )

    # ``--pack`` is an optional CLI shortcut for the spec-level "pack" field
    # (the engine resolves names/aliases like 技术风格 → dark_gold). CLI wins.
    if getattr(args, "pack", None):
        spec["pack"] = args.pack

    with tempfile.TemporaryDirectory(prefix="ppt_build_") as tmpdir, _shims.use_workdir(tmpdir):
        result = builder.build_from_spec(
            spec=spec,
            output_filename=out_path.name,
            engine=args.engine,
            style=args.style,
            theme=args.theme,
        )
        produced = Path(tmpdir) / out_path.name
        if not produced.is_file():
            _die("BuildFailed", f"engine reported success but {produced} is missing")
        shutil.move(str(produced), str(out_path))

    _emit(
        {
            "ok": True,
            "output": str(out_path),
            "meta": {
                "engine": result.get("engine"),
                "theme": result.get("theme"),
                "style": result.get("style"),
                "slide_count": result.get("slide_count"),
                "size_bytes": result.get("size_bytes"),
                "layout_warnings": result.get("layout_warnings", []),
                "font_warnings": result.get("font_warnings", []),
            },
        }
    )


def cmd_build_js(args: argparse.Namespace) -> None:
    """Run an agent-authored pptxgenjs Node script — the *freeform* build path.

    ``build`` feeds a JSON spec to the fixed template engine: safe, cheap,
    but bounded by the ~24 built-in layouts. ``build-js`` instead executes
    raw Node.js the agent wrote against the pptxgenjs API directly, so every
    slide can be laid out pixel-by-pixel — no layout is off-limits. The CLI
    only resolves pptxgenjs, runs the script under a timeout, and verifies a
    .pptx landed where expected.

    Contract: the script MUST write its deck to ``process.env.PPT_OUT_PATH``
    (e.g. ``await pres.writeFile({ fileName: process.env.PPT_OUT_PATH })``).
    """
    script = _ensure_input(args.script)
    out_path = _ensure_pptx_output(args.output)

    node = shutil.which("node") or shutil.which("nodejs")
    if not node:
        _die(
            "NodeNotFound",
            "Node.js not found in PATH — freeform build needs node + pptxgenjs. "
            "Fall back to `ppt-cli build` with a JSON spec.",
            field="script",
        )

    env = os.environ.copy()
    # A bare ``require("pptxgenjs")`` only resolves when the script sits
    # inside a node_modules tree. Freeform scripts live in the workspace, so
    # point NODE_PATH at the global install (pptxgenjs is installed `-g` in
    # the sandbox / mcp images — see Dockerfile.mcp).
    npm_root = ""
    try:
        npm_root = subprocess.run(
            ["npm", "root", "-g"], capture_output=True, text=True, timeout=10
        ).stdout.strip()
    except Exception:
        pass  # npm absent — rely on require()'s default resolution
    node_path = os.pathsep.join(p for p in (npm_root, env.get("NODE_PATH", "")) if p)
    if node_path:
        env["NODE_PATH"] = node_path
    env["PPT_OUT_PATH"] = str(out_path)

    try:
        proc = subprocess.run(
            [node, str(script)],
            cwd=str(script.parent),  # so relative image/asset paths resolve
            capture_output=True,
            text=True,
            timeout=args.timeout,
            env=env,
        )
    except subprocess.TimeoutExpired:
        _die("FreeformTimeout", f"freeform script did not finish within {args.timeout}s")

    if proc.returncode != 0:
        _die(
            "FreeformBuildFailed",
            proc.stderr.strip() or proc.stdout.strip() or f"node exited {proc.returncode}",
            field="script",
        )

    if not out_path.is_file():
        # The script ignored PPT_OUT_PATH and wrote elsewhere — accept it
        # only if exactly one .pptx appeared next to the script.
        candidates = sorted(script.parent.glob("*.pptx"))
        if len(candidates) == 1:
            shutil.move(str(candidates[0]), str(out_path))
        else:
            _die(
                "FreeformNoOutput",
                "script exited 0 but produced no .pptx at PPT_OUT_PATH — end your "
                "script with `await pres.writeFile({ fileName: process.env.PPT_OUT_PATH })`",
                field="script",
            )

    _emit(
        {
            "ok": True,
            "output": str(out_path),
            "meta": {
                "mode": "freeform",
                "size_bytes": out_path.stat().st_size,
                "node_stdout": (proc.stdout or "").strip()[-500:],
            },
        }
    )


def cmd_info(args: argparse.Namespace) -> None:
    """Slide count + per-slide titles — opens the deck once.

    The reader module exposes ``get_slide_count`` and ``get_slide_content``,
    but each call re-parses the .pptx. A 20-slide deck would unzip + parse
    21 times. We bypass them and use python-pptx directly so the file is
    opened once and we short-circuit each slide at the first non-empty
    paragraph (no full-text extraction needed for the title heuristic).
    """
    from pptx import Presentation

    src = _ensure_input(args.input)
    prs = Presentation(str(src))
    slides: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    title = line
                    break
            if title:
                break
        slides.append({"index": i, "title": title})
    _emit({"ok": True, "input": str(src), "slide_count": len(slides), "slides": slides})


def cmd_slide_count(args: argparse.Namespace) -> None:
    from engine import _shims, reader

    src = _ensure_input(args.input)
    with _shims.use_workdir(src.parent):
        info = reader.get_slide_count(input_filename=src.name)
    _emit({"ok": True, "input": str(src), "slide_count": info["slide_count"]})


def cmd_extract(args: argparse.Namespace) -> None:
    from engine import _shims, qa, reader

    src = _ensure_input(args.input)
    with _shims.use_workdir(src.parent):
        if args.slide is not None:
            sc = reader.get_slide_content(input_filename=src.name, slide_index=args.slide)
            _emit({"ok": True, "input": str(src), **sc})
        else:
            blocks = qa.extract_all_text(file_path=str(src))
            _emit(
                {
                    "ok": True,
                    "input": str(src),
                    "block_count": len(blocks),
                    "text_blocks": blocks,
                }
            )


def cmd_check_placeholders(args: argparse.Namespace) -> None:
    from engine import _shims, qa

    src = _ensure_input(args.input)
    patterns = args.patterns or None
    with _shims.use_workdir(src.parent):
        hits = qa.check_placeholders(file_path=str(src), patterns=patterns)
    _emit(
        {
            "ok": True,
            "input": str(src),
            "hit_count": len(hits),
            "is_clean": len(hits) == 0,
            "hits": hits,
        }
    )


def cmd_thumbnails(args: argparse.Namespace) -> None:
    from engine import _shims, thumbnails

    src = _ensure_input(args.input)
    out_dir = Path(args.output_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    if not (30 <= args.dpi <= 600):
        _die("ValueError", f"--dpi must be between 30 and 600, got {args.dpi}", field="dpi")
    if not (1 <= args.quality <= 100):
        _die("ValueError", f"--quality must be between 1 and 100, got {args.quality}", field="quality")

    with tempfile.TemporaryDirectory(prefix="ppt_thumb_") as tmpdir, _shims.use_workdir(tmpdir):
        # materialise input into workdir
        local_in = Path(tmpdir) / src.name
        shutil.copyfile(src, local_in)
        result = thumbnails.render_thumbnails(
            input_filename=src.name,
            output_prefix=args.prefix,
            dpi=args.dpi,
            quality=args.quality,
        )
        produced_paths: list[dict[str, Any]] = []
        for thumb in result.get("thumbnails", []):
            inner = Path(tmpdir) / thumb["filename"]
            final = out_dir / thumb["filename"]
            shutil.move(str(inner), str(final))
            produced_paths.append(
                {
                    "slide_index": thumb["slide_index"],
                    "path": str(final),
                    "size_bytes": thumb.get("size_bytes"),
                }
            )

    _emit(
        {
            "ok": True,
            "input": str(src),
            "output_dir": str(out_dir),
            "slide_count": result.get("slide_count"),
            "thumbnails": produced_paths,
            "dpi": args.dpi,
        }
    )


def _run_edit(
    fn,
    src: Path,
    out_path: Path,
    **kwargs,
) -> dict[str, Any]:
    """Common scaffold for single-slide editing ops in editor.py."""
    from engine import _shims

    with tempfile.TemporaryDirectory(prefix="ppt_edit_") as tmpdir, _shims.use_workdir(tmpdir):
        local_in = Path(tmpdir) / src.name
        shutil.copyfile(src, local_in)
        out_basename = out_path.name
        # editor functions expect input + output filenames in workdir
        result = fn(input_filename=src.name, output_filename=out_basename, **kwargs)
        produced = Path(tmpdir) / out_basename
        if not produced.is_file():
            _die("EditFailed", f"engine reported success but {produced} is missing")
        shutil.move(str(produced), str(out_path))
    return result


def _parse_content_for_slide_type(slide_type: str, raw: str | None) -> dict | None:
    """Map the CLI's free-text ``--content`` into the dict shape the engine wants.

    Builder.add_slide takes ``content: dict | None`` whose keys depend on the
    slide type (cover: subtitle/tagline/body; section: subtitle; content/summary:
    bullets). The CLI receives a flat string for ergonomics; this function does
    the per-type mapping so callers don't have to hand-craft JSON for the
    common cases.

    For content/summary slides we split the string by newlines into bullets —
    matching what the CLI's --content help text already promises. Anything
    further (icon_rows / stat_callout / etc.) needs the full ``build`` flow
    with a real spec, not this single-slide append.
    """
    if raw is None or raw == "":
        return None
    s = raw.strip()
    if slide_type in ("content", "summary"):
        bullets = [line.strip() for line in s.splitlines() if line.strip()]
        return {"bullets": bullets} if bullets else None
    if slide_type in ("cover",):
        return {"subtitle": s}
    if slide_type == "section":
        return {"subtitle": s}
    if slide_type == "toc":
        items = [line.strip() for line in s.splitlines() if line.strip()]
        return {"items": items} if items else None
    return None  # unreachable — slide_type already validated by caller


def cmd_add_slide(args: argparse.Namespace) -> None:
    from engine import _shims, builder

    src = _ensure_input(args.input)
    out_path = _ensure_pptx_output(args.output)
    # ``--type`` and ``--style`` are enforced by argparse ``choices=``.

    content_dict = _parse_content_for_slide_type(args.type, args.content)

    with tempfile.TemporaryDirectory(prefix="ppt_addslide_") as tmpdir, _shims.use_workdir(tmpdir):
        shutil.copyfile(src, Path(tmpdir) / src.name)
        result = builder.add_slide(
            input_filename=src.name,
            output_filename=out_path.name,
            slide_type=args.type,
            title=args.title,
            content=content_dict,
            theme=args.theme,
            style=args.style,
        )
        produced = Path(tmpdir) / out_path.name
        if not produced.is_file():
            _die("AddSlideFailed", f"engine reported success but {produced} is missing")
        shutil.move(str(produced), str(out_path))

    _emit({"ok": True, "output": str(out_path), "meta": result})


def cmd_set_title(args: argparse.Namespace) -> None:
    from engine import editor

    src = _ensure_input(args.input)
    out_path = _ensure_pptx_output(args.output)
    result = _run_edit(
        editor.set_slide_title,
        src,
        out_path,
        slide_index=args.slide,
        title=args.title,
    )
    _emit({"ok": True, "output": str(out_path), "meta": result})


def cmd_add_text(args: argparse.Namespace) -> None:
    from engine import editor

    src = _ensure_input(args.input)
    out_path = _ensure_pptx_output(args.output)
    position = _parse_position(args.position)
    result = _run_edit(
        editor.add_text_to_slide,
        src,
        out_path,
        slide_index=args.slide,
        text=args.text,
        position=position,
        font_size=args.font_size,
        color_hex=args.color,
        bold=args.bold,
    )
    _emit({"ok": True, "output": str(out_path), "meta": result})


def cmd_insert_image(args: argparse.Namespace) -> None:
    from engine import _shims, editor

    src = _ensure_input(args.input)
    img_src = _ensure_input(args.image)
    out_path = _ensure_pptx_output(args.output)
    position = _parse_position(args.position)
    with tempfile.TemporaryDirectory(prefix="ppt_img_") as tmpdir, _shims.use_workdir(tmpdir):
        shutil.copyfile(src, Path(tmpdir) / src.name)
        # Editor finds the image by bare filename, so stage it alongside the deck.
        img_local_name = img_src.name
        shutil.copyfile(img_src, Path(tmpdir) / img_local_name)
        result = editor.insert_image(
            input_filename=src.name,
            output_filename=out_path.name,
            slide_index=args.slide,
            image_filename=img_local_name,
            position=position,
        )
        produced = Path(tmpdir) / out_path.name
        if not produced.is_file():
            _die("InsertImageFailed", f"engine reported success but {produced} is missing")
        shutil.move(str(produced), str(out_path))
    _emit({"ok": True, "output": str(out_path), "meta": result})


def cmd_delete_slide(args: argparse.Namespace) -> None:
    from engine import editor

    src = _ensure_input(args.input)
    out_path = _ensure_pptx_output(args.output)
    result = _run_edit(editor.delete_slide, src, out_path, slide_index=args.slide)
    _emit({"ok": True, "output": str(out_path), "meta": result})


def cmd_to_pdf(args: argparse.Namespace) -> None:
    from engine import _shims

    src = _ensure_input(args.input)
    out_path = _ensure_pdf_output(args.output)
    with tempfile.TemporaryDirectory(prefix="ppt_pdf_") as tmpdir, _shims.use_workdir(tmpdir):
        shutil.copyfile(src, Path(tmpdir) / src.name)
        result = _shims.to_pdf(input_filename=src.name, output_filename=out_path.name)
        produced = Path(tmpdir) / out_path.name
        if not produced.is_file():
            _die("PDFExportFailed", f"engine reported success but {produced} is missing")
        shutil.move(str(produced), str(out_path))
    _emit({"ok": True, "output": str(out_path), "meta": result})


def cmd_list_themes(_: argparse.Namespace) -> None:
    _emit(
        {
            "ok": True,
            "palettes": sorted(list_palettes()),
            "aliases": list_aliases(),
        }
    )


def cmd_list_styles(_: argparse.Namespace) -> None:
    _emit({"ok": True, "styles": list_recipes()})


def cmd_list_slide_types(_: argparse.Namespace) -> None:
    _emit({"ok": True, "slide_types": list(_VALID_SLIDE_TYPES)})


# ── argument parser ─────────────────────────────────────────────────


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="ppt-cli",
        description="ppt-design skill CLI — build, inspect, edit, render, and export .pptx decks",
    )
    sub = p.add_subparsers(dest="cmd", required=True)

    # build
    sp = sub.add_parser("build", help="build a complete .pptx from a JSON spec (primary entry point)")
    sp.add_argument("--spec", required=True, help="path to .json spec file, or inline JSON string")
    sp.add_argument("--output", required=True, help="output .pptx path")
    sp.add_argument("--engine", default="pptxgenjs", choices=_VALID_ENGINES,
                    help="pptxgenjs (default) or python-pptx")
    sp.add_argument("--style", default="soft", choices=_VALID_STYLES,
                    help="default soft; sharp / rounded / pill also available")
    sp.add_argument("--theme", default=None,
                    help="palette name or alias; omit for default swiss_klein. "
                         "Run ``ppt-cli list-themes`` for the full catalog.")
    sp.add_argument("--pack", default=None,
                    help="design pack: dark_gold (技术风格/深蓝金) for the tech style; "
                         "omit for the default style. Also settable as \"pack\" in the spec JSON.")
    sp.set_defaults(func=cmd_build)

    # build-js (freeform — agent-authored pptxgenjs script)
    sp = sub.add_parser(
        "build-js",
        help="run an agent-authored pptxgenjs Node script — the freeform build path",
    )
    sp.add_argument("--script", required=True,
                    help="path to the .js file you wrote against the pptxgenjs API")
    sp.add_argument("--output", required=True, help="output .pptx path")
    sp.add_argument("--timeout", type=float, default=90.0,
                    help="seconds to wait for the script (default 90)")
    sp.set_defaults(func=cmd_build_js)

    # info
    sp = sub.add_parser("info", help="overview: slide count + per-slide titles")
    sp.add_argument("input", help="path to .pptx")
    sp.set_defaults(func=cmd_info)

    sp = sub.add_parser("slide-count", help="just the slide count (cheap probe)")
    sp.add_argument("input")
    sp.set_defaults(func=cmd_slide_count)

    # extract
    sp = sub.add_parser("extract", help="extract text — one slide if --slide, else every slide")
    sp.add_argument("input")
    sp.add_argument("--slide", type=int, default=None, help="0-based slide index (omit for all)")
    sp.set_defaults(func=cmd_extract)

    # check-placeholders
    sp = sub.add_parser("check-placeholders", help="scan for placeholder strings (xxxx / lorem / 占位 …)")
    sp.add_argument("input")
    sp.add_argument("--patterns", nargs="*", default=None, help="override patterns; omit for default set")
    sp.set_defaults(func=cmd_check_placeholders)

    # thumbnails
    sp = sub.add_parser("thumbnails", help="render each slide as a JPG into --output-dir")
    sp.add_argument("input")
    sp.add_argument("--output-dir", required=True, help="directory to write the JPGs into")
    sp.add_argument("--prefix", default="slide", help="filename prefix (default: slide)")
    sp.add_argument("--dpi", type=int, default=120, help="30–600 (default 120)")
    sp.add_argument("--quality", type=int, default=85, help="1–100 (default 85)")
    sp.set_defaults(func=cmd_thumbnails)

    # add-slide (single-slide append on an existing deck)
    sp = sub.add_parser("add-slide", help="append a typed slide to an existing deck")
    sp.add_argument("input")
    sp.add_argument("--output", required=True)
    sp.add_argument("--type", required=True, choices=_VALID_SLIDE_TYPES,
                    help="slide type for the appended page")
    sp.add_argument("--title", default=None)
    sp.add_argument(
        "--content",
        default=None,
        help=(
            "for cover/section: subtitle text; for content/summary: bullets, one per line "
            "(split on \\n internally); for toc: items, one per line. "
            "More complex content (icon_rows / stat_callout / ...) needs a full `build` with a spec."
        ),
    )
    sp.add_argument("--theme", default="default",
                    help="palette name or alias; ``default`` ≡ swiss_klein. "
                         "Run ``ppt-cli list-themes`` for the full catalog.")
    sp.add_argument("--style", default="soft", choices=_VALID_STYLES,
                    help="default soft; sharp / rounded / pill also available")
    sp.set_defaults(func=cmd_add_slide)

    # set-title
    sp = sub.add_parser("set-title", help="replace a slide's title text")
    sp.add_argument("input")
    sp.add_argument("--output", required=True)
    sp.add_argument("--slide", type=int, required=True, help="0-based slide index")
    sp.add_argument("--title", required=True)
    sp.set_defaults(func=cmd_set_title)

    # add-text
    sp = sub.add_parser("add-text", help="insert a free-form text box on a slide")
    sp.add_argument("input")
    sp.add_argument("--output", required=True)
    sp.add_argument("--slide", type=int, required=True)
    sp.add_argument("--text", required=True)
    sp.add_argument("--position", default=None, help="comma-separated inches: 'L,T' or 'L,T,W,H'")
    sp.add_argument("--font-size", type=int, default=14, help="point size, default 14")
    sp.add_argument("--color", default=None, help="6-char hex without #, e.g. 1A2B3C")
    sp.add_argument("--bold", action="store_true")
    sp.set_defaults(func=cmd_add_text)

    # insert-image
    sp = sub.add_parser("insert-image", help="insert an image on a slide")
    sp.add_argument("input")
    sp.add_argument("--output", required=True)
    sp.add_argument("--slide", type=int, required=True)
    sp.add_argument("--image", required=True, help="path to image file (.png/.jpg/.svg etc.)")
    sp.add_argument("--position", default=None, help="comma-separated inches: 'L,T' or 'L,T,W,H'")
    sp.set_defaults(func=cmd_insert_image)

    # delete-slide
    sp = sub.add_parser("delete-slide", help="remove a slide by index")
    sp.add_argument("input")
    sp.add_argument("--output", required=True)
    sp.add_argument("--slide", type=int, required=True)
    sp.set_defaults(func=cmd_delete_slide)

    # to-pdf
    sp = sub.add_parser("to-pdf", help="convert a .pptx to a .pdf via LibreOffice headless")
    sp.add_argument("input")
    sp.add_argument("--output", required=True)
    sp.set_defaults(func=cmd_to_pdf)

    # introspection helpers
    sp = sub.add_parser("list-themes", help="list all palettes + aliases")
    sp.set_defaults(func=cmd_list_themes)

    sp = sub.add_parser("list-styles", help="list visual styles (sharp/soft/rounded/pill)")
    sp.set_defaults(func=cmd_list_styles)

    sp = sub.add_parser("list-slide-types", help="list valid slide types for add-slide")
    sp.set_defaults(func=cmd_list_slide_types)

    return p


def main(argv: list[str] | None = None) -> int:
    parser = _build_parser()
    args = parser.parse_args(argv)
    try:
        args.func(args)
    except FileNotFoundError as exc:
        _die("FileNotFound", str(exc))
    except ValueError as exc:
        _die("ValueError", str(exc))
    except Exception as exc:  # pragma: no cover — last-resort surface
        _die(type(exc).__name__, str(exc))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
