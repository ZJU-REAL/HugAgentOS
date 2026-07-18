#!/usr/bin/env python3
"""
文件级脱敏模块 - 支持 Word、PowerPoint、PDF、图片、HTML、Markdown 等格式。

与 mask.py / dataframe_mask.py 配合，提供完整的多格式脱敏能力。

混合策略：
  - 原生支持（内置解析 + 脱敏 + 同格式输出）：DOCX / PPTX / XLSX / CSV / TXT / MD
  - 原生支持（提取文本 + 脱敏 → .txt 输出）：HTML / PDF / 图片(PNG/JPG/...)
  - 未知格式回退：尝试读取为文本 → 脱敏 → .txt 输出
  - 扫描版 PDF / 复杂格式：委托专项 skill 提取文本 → mask.py 脱敏

使用示例:
  # 命令行（-o 可选，默认：原文件名_mask.原后缀）
  python file_mask.py -i report.docx                    # → report_mask.docx
  python file_mask.py -i slides.pptx --strategy blur    # → slides_mask.pptx
  python file_mask.py -i document.pdf                   # → document_mask.pdf
  python file_mask.py -i page.html                      # → page_mask.txt
  python file_mask.py -i notes.md                       # → notes_mask.md
  python file_mask.py -i screenshot.png --ocr-lang chi_sim+eng  # → screenshot_mask.txt

  # 编程调用（output_path 可选，传 None 自动生成）
  from file_mask import mask_document
  mask_document('report.docx', None, strategy='hash')   # → report_mask.docx
  mask_document('report.docx', 'custom.docx', strategy='hash')  # → custom.docx
"""

import argparse
import json
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from mask import mask_text


# ── 格式检测 ─────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    '.docx': 'docx',
    '.pptx': 'pptx',
    '.pdf': 'pdf',
    '.txt': 'txt',
    '.xlsx': 'xlsx',
    '.xls': 'xlsx',
    '.csv': 'csv',
    '.html': 'html',
    '.htm': 'html',
    '.md': 'md',
    '.markdown': 'md',
    # 图片格式 → OCR 提取文字
    '.png': 'image',
    '.jpg': 'image',
    '.jpeg': 'image',
    '.gif': 'image',
    '.bmp': 'image',
    '.webp': 'image',
    '.tiff': 'image',
    '.tif': 'image',
}


def detect_format(filepath):
    """根据扩展名检测文件格式"""
    ext = os.path.splitext(filepath)[1].lower()
    return SUPPORTED_EXTENSIONS.get(ext)


def generate_default_output(input_path, fmt=None):
    """
    根据输入文件路径和格式生成默认输出路径。
    命名规范：原文件名_mask.原后缀
      例：report.docx → report_mask.docx
         data.xlsx  → data_mask.xlsx
         page.html  → page_mask.txt（HTML 输出固定 .txt）
         img.png    → img_mask.txt（图片 OCR 输出 .txt）
    """
    dir_name = os.path.dirname(input_path) or '.'
    base = os.path.splitext(os.path.basename(input_path))[0]
    ext = os.path.splitext(input_path)[1].lower()

    # 以下格式固定输出为 .txt
    if fmt in ('html', 'image') or fmt is None:
        return os.path.join(dir_name, f'{base}_mask.txt')

    return os.path.join(dir_name, f'{base}_mask{ext}')


# ── DOCX 处理 ────────────────────────────────────────────────

def _iter_docx_paragraphs(doc):
    """遍历文档所有段落（含表格、页眉页脚）"""
    # 正文段落
    yield from doc.paragraphs

    # 表格中的段落
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs

    # 页眉页脚
    for section in doc.sections:
        for para in section.header.paragraphs:
            yield para
        for para in section.footer.paragraphs:
            yield para
        # 首页页眉页脚
        if section.different_first_page_header_footer:
            for para in section.first_page_header.paragraphs:
                yield para
            for para in section.first_page_footer.paragraphs:
                yield para


def _mask_docx_text(doc, strategy, seed, skip_phone, mapping=None):
    """对 DOCX 文档内容进行脱敏，保留原有格式。返回共享映射表。"""
    if mapping is None:
        mapping = {}
    for para in _iter_docx_paragraphs(doc):
        if not para.runs:
            continue

        # 收集所有 run 的文本和格式
        runs_data = [(run, run.text) for run in para.runs if run.text]

        if not runs_data:
            continue

        # 拼接全文 → 脱敏 → 按位置回写
        full_text = ''.join(t for _, t in runs_data)
        masked_text, _ = mask_text(
            full_text, strategy=strategy, mapping=mapping,
            seed=seed, skip_phone=skip_phone
        )

        if masked_text == full_text:
            continue  # 没有变化，无需修改

        # 清空所有 run，写入脱敏后的文本到第一个 run，其余清空
        for i, (run, _) in enumerate(runs_data):
            if i == 0:
                run.text = masked_text
            else:
                run.text = ''

    return mapping


def mask_docx(input_path, output_path, strategy='hash', seed=42,
              skip_phone=False, map_path=None):
    """
    对 Word(.docx) 文档进行脱敏。

    参数:
        input_path: 输入 .docx 文件路径
        output_path: 输出 .docx 文件路径
        strategy: 脱敏策略
        seed: 随机种子
        skip_phone: 是否跳过电话号码
        map_path: 映射表输出路径

    返回:
        映射字典
    """
    try:
        from docx import Document
    except ImportError:
        print(
            '错误：处理 Word 文档需要安装 python-docx 库。\n'
            '  请执行：pip install python-docx',
            file=sys.stderr
        )
        sys.exit(1)

    try:
        doc = Document(input_path)
    except Exception as e:
        print(
            f'错误：无法打开 Word 文档 "{input_path}" —— {e}\n'
            f'  请确认：\n'
            f'  1. 文件是有效的 .docx 格式（不是旧版 .doc）\n'
            f'  2. 文件未损坏\n'
            f'  3. 文件未被其他程序占用',
            file=sys.stderr
        )
        sys.exit(1)

    mapping = _mask_docx_text(doc, strategy, seed, skip_phone)

    try:
        doc.save(output_path)
    except PermissionError:
        print(
            f'错误：无法写入输出文件 "{output_path}" —— 权限不足。\n'
            f'  请确认：\n'
            f'  1. 输出目录有写入权限\n'
            f'  2. 同名文件未被其他程序（如 Word）打开',
            file=sys.stderr
        )
        sys.exit(1)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    return mapping


# ── PPTX 处理 ────────────────────────────────────────────────

def _iter_pptx_text_frames(prs):
    """遍历 PPTX 中所有文本框架"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                yield shape.text_frame
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
            # 组合形状递归
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                yield from _iter_group_shapes(shape, set())


def _iter_group_shapes(group_shape, visited):
    """递归遍历组合形状"""
    gid = id(group_shape)
    if gid in visited:
        return
    visited.add(gid)
    try:
        for shape in group_shape.shapes:
            if shape.has_text_frame:
                yield shape.text_frame
            if hasattr(shape, 'has_table') and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
            if hasattr(shape, 'shape_type') and shape.shape_type == 6:
                yield from _iter_group_shapes(shape, visited)
    except Exception:
        pass


def _mask_pptx_text(prs, strategy, seed, skip_phone, mapping=None):
    """对 PPTX 中所有文本框架进行脱敏。返回共享映射表。"""
    if mapping is None:
        mapping = {}
    for tf in _iter_pptx_text_frames(prs):
        for para in tf.paragraphs:
            if not para.runs:
                continue
            runs_data = [(run, run.text) for run in para.runs if run.text]
            if not runs_data:
                continue

            full_text = ''.join(t for _, t in runs_data)
            masked_text, _ = mask_text(
                full_text, strategy=strategy, mapping=mapping,
                seed=seed, skip_phone=skip_phone
            )

            if masked_text == full_text:
                continue

            for i, (run, _) in enumerate(runs_data):
                if i == 0:
                    run.text = masked_text
                else:
                    run.text = ''

    return mapping


def mask_pptx(input_path, output_path, strategy='hash', seed=42,
              skip_phone=False, map_path=None):
    """
    对 PowerPoint(.pptx) 演示文稿进行脱敏。

    参数:
        input_path: 输入 .pptx 文件路径
        output_path: 输出 .pptx 文件路径
        strategy: 脱敏策略
        seed: 随机种子
        skip_phone: 是否跳过电话号码
        map_path: 映射表输出路径

    返回:
        映射字典
    """
    try:
        from pptx import Presentation
    except ImportError:
        print(
            '错误：处理 PPT 需要安装 python-pptx 库。\n'
            '  请执行：pip install python-pptx',
            file=sys.stderr
        )
        sys.exit(1)

    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(
            f'错误：无法打开 PPT 文件 "{input_path}" —— {e}\n'
            f'  请确认：\n'
            f'  1. 文件是有效的 .pptx 格式\n'
            f'  2. 文件未损坏\n'
            f'  3. 文件未被其他程序占用',
            file=sys.stderr
        )
        sys.exit(1)

    mapping = _mask_pptx_text(prs, strategy, seed, skip_phone)

    try:
        prs.save(output_path)
    except PermissionError:
        print(
            f'错误：无法写入输出文件 "{output_path}" —— 权限不足。\n'
            f'  请确认输出目录有写入权限，且同名文件未被其他程序打开。',
            file=sys.stderr
        )
        sys.exit(1)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    return mapping


# ── PDF 处理 ─────────────────────────────────────────────────

def _find_cjk_font():
    """查找系统可用的中文字体"""
    import platform
    font_paths = []

    if platform.system() == 'Windows':
        font_dir = os.path.join(os.environ.get('WINDIR', 'C:\\Windows'), 'Fonts')
        candidates = [
            'msyh.ttc', 'msyh.ttf',     # 微软雅黑
            'simsun.ttc', 'simsun.ttf',  # 宋体
            'simhei.ttf',                # 黑体
            'kaiu.ttf',                  # 楷体
            'Deng.ttf', 'Dengb.ttf',    # 等线
        ]
        for name in candidates:
            path = os.path.join(font_dir, name)
            if os.path.exists(path):
                font_paths.append(path)
    elif platform.system() == 'Darwin':
        font_dir = '/System/Library/Fonts'
        candidates = [
            'PingFang.ttc', 'STHeiti Light.ttc',
            'Hiragino Sans GB.ttc', 'AppleSDGothicNeo.ttc',
        ]
        for name in candidates:
            path = os.path.join(font_dir, name)
            if os.path.exists(path):
                font_paths.append(path)
    else:  # Linux
        font_dirs = ['/usr/share/fonts', '/usr/local/share/fonts']
        for fd in font_dirs:
            for root, _, files in os.walk(fd):
                for f in files:
                    if any(name in f.lower() for name in
                           ['wqy', 'noto', 'cjk', 'droid', 'wenquan']):
                        font_paths.append(os.path.join(root, f))

    return font_paths[0] if font_paths else None


def mask_pdf(input_path, output_path, strategy='hash', seed=42,
             skip_phone=False, map_path=None, as_text=False):
    """
    对 PDF 文档进行脱敏。

    输出策略:
      - 默认: 生成新的 PDF，尽可能包含脱敏后的文本
      - as_text=True: 生成 .txt 文本文件

    参数:
        input_path: 输入 .pdf 文件路径
        output_path: 输出文件路径
        strategy: 脱敏策略
        seed: 随机种子
        skip_phone: 是否跳过电话号码
        map_path: 映射表输出路径
        as_text: 是否输出为纯文本（默认尝试生成 PDF）

    返回:
        映射字典
    """
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        print('错误：需要安装 PyPDF2: pip install PyPDF2', file=sys.stderr)
        sys.exit(1)

    # 提取 PDF 文本
    reader = PdfReader(input_path)
    pages_text = []
    for page in reader.pages:
        text = page.extract_text()
        pages_text.append(text if text else '')

    # 逐页脱敏，使用全局映射表保证跨页一致性
    mapping = {}
    masked_pages = []
    for page_text in pages_text:
        if page_text.strip():
            masked, _ = mask_text(
                page_text, strategy=strategy, mapping=mapping,
                seed=seed, skip_phone=skip_phone
            )
            masked_pages.append(masked)
        else:
            masked_pages.append('')

    masked_text = '\n'.join(masked_pages)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    # 输出
    if as_text or output_path.lower().endswith('.txt'):
        output = output_path if output_path.lower().endswith('.txt') else output_path + '.txt'
        with open(output, 'w', encoding='utf-8') as f:
            f.write(masked_text)
        return mapping

    # 尝试生成 PDF
    try:
        _write_masked_pdf(masked_pages, output_path)
    except Exception as e:
        # 回退到文本输出
        print(f'PDF 生成失败 ({e})，保存为文本文件', file=sys.stderr)
        txt_path = os.path.splitext(output_path)[0] + '.txt'
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(masked_text)
        print(f'文本输出: {txt_path}', file=sys.stderr)

    return mapping


def _write_masked_pdf(masked_pages, output_path):
    """生成包含脱敏文本的新 PDF，每页已独立脱敏，精确分页。"""
    try:
        from fpdf import FPDF
    except ImportError:
        raise RuntimeError('需要安装 fpdf2: pip install fpdf2')

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # 查找中文字体
    cjk_font = _find_cjk_font()
    if cjk_font:
        pdf.add_font('CJK', '', cjk_font, uni=True)
        pdf.add_font('CJK', 'B', cjk_font, uni=True)
        font_name = 'CJK'
    else:
        font_name = 'Helvetica'

    for page_text in masked_pages:
        pdf.add_page()
        pdf.set_font(font_name, '', 10)

        if not page_text.strip():
            continue

        lines = page_text.split('\n')
        for line in lines:
            if not line.strip():
                pdf.cell(0, 5, '', ln=True)
                continue
            while len(line) > 100:
                pdf.cell(0, 5, line[:100], ln=True)
                line = line[100:]
            pdf.cell(0, 5, line, ln=True)

    pdf.output(output_path)


# ── HTML 处理 ────────────────────────────────────────────────

def _extract_html_text(filepath):
    """从 HTML 文件中提取可见文本，保留段落结构。使用 stdlib，无外部依赖。"""
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self.skip_stack = 0       # 嵌套跳过层数
            self.skip_tags = {'script', 'style', 'noscript', 'iframe', 'svg', 'code'}
            self.block_tags = {
                'div', 'p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
                'li', 'tr', 'br', 'hr', 'section', 'article',
                'header', 'footer', 'nav', 'main', 'aside',
                'pre', 'blockquote', 'table', 'ul', 'ol', 'dl',
                'form', 'fieldset', 'figcaption', 'figure',
            }
            self.heading_tags = {'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}
            self._last_was_newline = False

        def handle_starttag(self, tag, attrs):
            tag_lower = tag.lower()
            if tag_lower in self.skip_tags:
                self.skip_stack += 1
            if tag_lower in self.heading_tags:
                self._emit_newline()

        def handle_endtag(self, tag):
            tag_lower = tag.lower()
            if tag_lower in self.skip_tags:
                self.skip_stack = max(0, self.skip_stack - 1)
                return
            if tag_lower in self.block_tags:
                self._emit_newline()

        def handle_data(self, data):
            if self.skip_stack > 0:
                return
            text = data.strip()
            if text:
                self.parts.append(text)
                self.parts.append(' ')
                self._last_was_newline = False

        def _emit_newline(self):
            if not self._last_was_newline:
                self.parts.append('\n')
                self._last_was_newline = True

        def get_text(self):
            return ''.join(self.parts).strip()

    extractor = TextExtractor()
    # 尝试多种编码读取 HTML 文件
    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(filepath, 'r', encoding=encoding) as f:
                extractor.feed(f.read())
            break
        except (UnicodeDecodeError, UnicodeError):
            extractor = TextExtractor()  # 重置解析器
            continue
    else:
        # 最后仍用 utf-8 尝试
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            extractor.feed(f.read())
    return extractor.get_text()


def mask_html(input_path, output_path, strategy='hash', seed=42,
              skip_phone=False, map_path=None):
    """
    对 HTML 文件进行脱敏。
    提取可见文本内容 → 脱敏 → 输出为结构化 .txt 文本。

    注意：HTML 脱敏不支持往返重建（无法写回原 HTML 结构），
          输出固定为 .txt 格式。
    """
    text = _extract_html_text(input_path)
    if not text.strip():
        print('警告: HTML 文件中未提取到可见文本内容', file=sys.stderr)
        return {}

    masked_text, mapping = mask_text(
        text, strategy=strategy, seed=seed, skip_phone=skip_phone
    )

    # HTML 只能输出为 .txt
    if not output_path.lower().endswith('.txt'):
        output_path = os.path.splitext(output_path)[0] + '.txt'

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(masked_text)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    return mapping


# ── 图片 OCR 处理 ────────────────────────────────────────────

# pytesseract 需要系统安装 Tesseract OCR 引擎，详见：
#   Windows: https://github.com/UB-Mannheim/tesseract/wiki
#   macOS:   brew install tesseract tesseract-lang
#   Linux:   apt-get install tesseract-ocr tesseract-ocr-chi-sim

def _ocr_image(filepath, lang='chi_sim+eng'):
    """
    使用 Tesseract OCR 从图片中提取文字。

    参数:
        filepath: 图片文件路径
        lang: OCR 语言代码，如 'chi_sim+eng'（中英混合）

    返回:
        提取到的文本字符串
    """
    try:
        from PIL import Image
    except ImportError:
        print('错误：需要安装 Pillow: pip install Pillow', file=sys.stderr)
        sys.exit(1)

    try:
        import pytesseract
    except ImportError:
        print('错误：需要安装 pytesseract: pip install pytesseract', file=sys.stderr)
        print('同时需要安装 Tesseract OCR 引擎：', file=sys.stderr)
        print('  Windows: https://github.com/UB-Mannheim/tesseract/wiki', file=sys.stderr)
        print('  macOS:   brew install tesseract tesseract-lang', file=sys.stderr)
        print('  Linux:   apt-get install tesseract-ocr tesseract-ocr-chi-sim', file=sys.stderr)
        sys.exit(1)

    # 尝试找到 Tesseract 安装路径（Windows 常见路径）
    if sys.platform == 'win32':
        common_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            os.path.expandvars(r'%LOCALAPPDATA%\Tesseract-OCR\tesseract.exe'),
        ]
        for p in common_paths:
            if os.path.exists(p):
                pytesseract.pytesseract.tesseract_cmd = p
                break

    img = Image.open(filepath)

    # WebP 需要转换（Tesseract 不支持直接读取 WebP）
    if filepath.lower().endswith('.webp'):
        img = img.convert('RGB')

    try:
        text = pytesseract.image_to_string(img, lang=lang)
    except pytesseract.TesseractError as e:
        # 语言包缺失？回退到默认英文
        print(f'OCR 警告 ({e})，回退到英文识别', file=sys.stderr)
        text = pytesseract.image_to_string(img, lang='eng')
    except pytesseract.TesseractNotFoundError:
        print('错误：未找到 Tesseract OCR 引擎！', file=sys.stderr)
        print('请先安装 Tesseract OCR：', file=sys.stderr)
        print('  Windows: https://github.com/UB-Mannheim/tesseract/wiki', file=sys.stderr)
        print('  macOS:   brew install tesseract tesseract-lang', file=sys.stderr)
        print('  Linux:   apt-get install tesseract-ocr tesseract-ocr-chi-sim', file=sys.stderr)
        sys.exit(1)

    return text.strip()


def mask_image(input_path, output_path, strategy='hash', seed=42,
               skip_phone=False, map_path=None, ocr_lang='chi_sim+eng'):
    """
    对图片文件进行脱敏：OCR 提取文字 → 脱敏 → 输出为 .txt。

    参数:
        input_path: 图片文件路径 (.png/.jpg/.gif/.bmp/.webp/.tiff)
        output_path: 输出文件路径（自动转为 .txt）
        strategy: 脱敏策略
        seed: 随机种子
        skip_phone: 是否跳过电话号码
        map_path: 映射表输出路径
        ocr_lang: OCR 语言（默认 chi_sim+eng 中英混合）

    返回:
        映射字典。如果 OCR 未提取到文字，返回空映射并输出空文件。
    """
    print(f'正在进行 OCR 文字识别 (语言: {ocr_lang})...', file=sys.stderr)

    text = _ocr_image(input_path, lang=ocr_lang)

    if not text:
        print('警告: OCR 未识别到文字内容，输出为空文件', file=sys.stderr)
        # 仍输出空文件，保持行为一致
        if not output_path.lower().endswith('.txt'):
            output_path = os.path.splitext(output_path)[0] + '.txt'
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('(OCR 未识别到文字)')
        return {}

    masked_text, mapping = mask_text(
        text, strategy=strategy, seed=seed, skip_phone=skip_phone
    )

    # 图片输出为 .txt
    if not output_path.lower().endswith('.txt'):
        output_path = os.path.splitext(output_path)[0] + '.txt'

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(masked_text)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    return mapping


# ── 纯文本 / Markdown 处理 ───────────────────────────────────

def mask_txt(input_path, output_path, strategy='hash', seed=42,
             skip_phone=False, map_path=None):
    """
    对纯文本文件进行脱敏。
    同时适用于 Markdown (.md) —— 数字被替换后 Markdown 语法不受影响。
    自动尝试多种编码读取。
    """
    from mask import mask_text

    # 多编码读取
    text = None
    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(input_path, 'r', encoding=encoding) as f:
                text = f.read()
            if encoding != 'utf-8':
                print(f'提示: 文件编码为 {encoding}，已自动适配', file=sys.stderr)
            break
        except (UnicodeDecodeError, UnicodeError):
            continue

    if text is None:
        raise ValueError(
            f'无法以文本编码读取 "{input_path}"。\n'
            f'  请确认该文件是否为文本文件，或尝试先转为 UTF-8 编码。'
        )

    mapping = {}
    placeholder_counter = [0]
    masked_text, mapping = mask_text(
        text, strategy=strategy, mapping=mapping,
        placeholder_counter=placeholder_counter,
        seed=seed, skip_phone=skip_phone
    )

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(masked_text)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    return mapping


# ── 未知格式回退 ─────────────────────────────────────────────

def mask_fallback(input_path, output_path, strategy='hash', seed=42,
                  skip_phone=False, map_path=None):
    """
    对未知格式文件尝试通用脱敏：
    1. 尝试以多编码文本读取
    2. 脱敏
    3. 输出为 .txt

    如果文件无法作为文本读取（二进制格式），给出明确的格式转换建议。
    """
    ext = os.path.splitext(input_path)[1].lower()

    # 尝试多种编码读取
    used_encoding = None
    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(input_path, 'r', encoding=encoding) as f:
                text = f.read()
            used_encoding = encoding
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    else:
        raise ValueError(
            f'无法以文本格式读取 "{input_path}"。\n'
            f'  该文件可能是二进制格式（{ext}），无法直接作为文本脱敏。\n'
            f'  建议：\n'
            f'    - 如果是扫描版 PDF → 用 PDF skill 做 OCR 提取文本，再用 mask.py 脱敏\n'
            f'    - 如果是图片            → 确认格式为 .png/.jpg 等常见图片格式\n'
            f'    - 如果是其他二进制格式   → 导出为 .txt 后再脱敏\n'
            f'    - 如果是文本格式但编码特殊 → 先手动转为 UTF-8 编码再脱敏'
        )

    if used_encoding and used_encoding != 'utf-8':
        print(f'提示: 文件编码为 {used_encoding}，已自动适配读取', file=sys.stderr)

    masked_text, mapping = mask_text(
        text, strategy=strategy, seed=seed, skip_phone=skip_phone
    )

    # 强制输出为 .txt
    if not output_path.lower().endswith('.txt'):
        output_path = os.path.splitext(output_path)[0] + '.txt'

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(masked_text)

    if map_path:
        with open(map_path, 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False, indent=2)

    return mapping


# ── 统一入口 ─────────────────────────────────────────────────

FORMAT_HANDLERS = {
    'docx': mask_docx,
    'pptx': mask_pptx,
    'pdf':  mask_pdf,
    'txt':  mask_txt,
    'html': mask_html,
    'md':   mask_txt,   # Markdown 本质是纯文本，数字替换不影响语法
    'image': mask_image,
}


def mask_document(input_path, output_path, strategy='hash', seed=42,
                  skip_phone=False, map_path=None, **kwargs):
    """
    统一的文档脱敏入口，自动识别格式并调用对应处理器。

    支持的格式：
      - 原生同格式输出：.docx / .pptx / .xlsx / .csv / .md / .txt
      - 提取文本 → .txt 输出：.html / .htm
      - 提取文本 → .pdf 或 .txt 输出：.pdf

    参数:
        input_path: 输入文件路径
        output_path: 输出文件路径（可选，默认自动生成：原文件名_mask.原后缀）
        strategy: 脱敏策略 (hash|placeholder|blur|range)
        seed: 随机种子
        skip_phone: 是否跳过电话号码
        map_path: 映射表输出路径
        **kwargs: 传递给特定处理器的额外参数

    返回:
        (格式类型, 映射字典)
    """
    fmt = detect_format(input_path)
    if output_path is None:
        output_path = generate_default_output(input_path, fmt=fmt)
        print(f'输出文件未指定，自动生成: {output_path}', file=sys.stderr)
    if fmt is None:
        # 未知格式 → 尝试通用文本脱敏回退
        print(f'注意: 未识别的格式，尝试作为文本读取并脱敏...', file=sys.stderr)
        mapping = mask_fallback(
            input_path, output_path,
            strategy=strategy, seed=seed,
            skip_phone=skip_phone, map_path=map_path
        )
        return 'fallback', mapping

    if fmt == 'xlsx' or fmt == 'csv':
        # 结构化数据委托给 dataframe_mask.py
        from dataframe_mask import mask_excel, mask_csv
        if fmt == 'xlsx':
            _, mapping = mask_excel(
                input_path, output_path, strategy=strategy,
                seed=seed, skip_phone=skip_phone, map_path=map_path, **kwargs
            )
        else:
            _, mapping = mask_csv(
                input_path, output_path, strategy=strategy,
                seed=seed, skip_phone=skip_phone, map_path=map_path, **kwargs
            )
        return fmt, mapping

    handler = FORMAT_HANDLERS[fmt]
    # 按格式过滤 kwargs：
    #   pdf   → as_text
    #   image → ocr_lang
    if fmt == 'pdf':
        filtered_kwargs = {k: v for k, v in kwargs.items() if k in ('as_text',)}
    elif fmt == 'image':
        filtered_kwargs = {k: v for k, v in kwargs.items() if k in ('ocr_lang',)}
    else:
        filtered_kwargs = {}
    mapping = handler(input_path, output_path, strategy=strategy,
                      seed=seed, skip_phone=skip_phone,
                      map_path=map_path, **filtered_kwargs)
    return fmt, mapping


# ── CLI ──────────────────────────────────────────────────────

def _validate_input_file(path):
    """验证输入文件：存在性、可读性、大小合理性。返回 (ok, message)"""
    if not os.path.exists(path):
        return False, f'错误：文件不存在 —— "{path}"\n   请确认文件路径是否正确，路径中的文件夹和文件名是否拼写无误。'
    if not os.path.isfile(path):
        return False, f'错误：路径不是文件 —— "{path}"\n   输入的路径是一个目录而非文件。'
    if not os.access(path, os.R_OK):
        return False, f'错误：文件无读取权限 —— "{path}"\n   请检查文件权限设置。'

    # 文件大小检查（> 500MB 提示）
    size_mb = os.path.getsize(path) / (1024 * 1024)
    if size_mb > 500:
        print(
            f'⚠ 警告：文件较大 ({size_mb:.0f} MB)，处理可能需要较长时间，请耐心等待...',
            file=sys.stderr
        )
    elif size_mb > 100:
        print(
            f'提示：文件 {size_mb:.0f} MB，处理中...',
            file=sys.stderr
        )
    return True, None


def main():
    parser = argparse.ArgumentParser(
        description='文档脱敏工具 - 支持 Word/PPT/PDF/图片/HTML/Markdown/TXT 格式',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python file_mask.py -i report.docx                    # 自动生成 report_mask.docx
  python file_mask.py -i report.docx -o masked.docx     # 手动指定输出
  python file_mask.py -i slides.pptx -s blur            # 自动生成 slides_mask.pptx
  python file_mask.py -i doc.pdf -m mapping.json        # 自动生成 doc_mask.pdf + 映射表
  python file_mask.py -i doc.pdf --as-text              # 输出为 doc_mask.txt
  python file_mask.py -i page.html                      # 自动生成 page_mask.txt
  python file_mask.py -i notes.md                       # 自动生成 notes_mask.md
  python file_mask.py -i screenshot.png --ocr-lang chi_sim+eng  # 生成 screenshot_mask.txt
  python file_mask.py -i unknown.json                   # 自动生成 unknown_mask.txt

批量处理示例:
  # 批量脱敏文件夹下所有 .docx（自动生成 *_mask.docx）
  for f in *.docx; do python file_mask.py -i "$f" -s hash; done
        """
    )
    parser.add_argument('-i', '--input', required=True, help='输入文件路径')
    parser.add_argument('-o', '--output', default=None,
                        help='输出文件路径（可选，默认：原文件名_mask.原后缀）')
    parser.add_argument('-s', '--strategy', default='hash',
                        choices=['hash', 'placeholder', 'blur', 'range'],
                        help='脱敏策略 (默认: hash)')
    parser.add_argument('--seed', type=int, default=42, help='随机种子 (默认: 42)')
    parser.add_argument('--skip-phone', action='store_true', help='跳过电话号码')
    parser.add_argument('-m', '--map-file', help='映射表输出路径 (JSON)')
    parser.add_argument('--as-text', action='store_true',
                        help='PDF 输出为文本而非 PDF')
    parser.add_argument('--ocr-lang', default='chi_sim+eng',
                        help='图片 OCR 语言 (默认: chi_sim+eng，纯英文用 eng，纯中文用 chi_sim)')

    args = parser.parse_args()

    # 1. 输入文件验证
    ok, err_msg = _validate_input_file(args.input)
    if not ok:
        print(err_msg, file=sys.stderr)
        sys.exit(1)

    # 2. 格式可识别性检查
    fmt = detect_format(args.input)
    if fmt is None:
        ext = os.path.splitext(args.input)[1].lower()
        print(
            f'提示：未识别的文件格式 ({ext})，将尝试作为文本读取并脱敏。',
            file=sys.stderr
        )
        print(
            f'  如果脱敏结果不正确，建议手动将文件内容导出为 .txt 后再脱敏。',
            file=sys.stderr
        )

    # 3. 自动生成输出路径（如未指定）
    if args.output is None:
        args.output = generate_default_output(args.input, fmt=fmt)
        print(f'输出文件未指定，自动生成: {args.output}', file=sys.stderr)

    # 4. 输出目录存在性检查
    out_dir = os.path.dirname(os.path.abspath(args.output))
    if out_dir and not os.path.exists(out_dir):
        try:
            os.makedirs(out_dir, exist_ok=True)
            print(f'已自动创建输出目录: {out_dir}', file=sys.stderr)
        except OSError as e:
            print(f'错误：无法创建输出目录 "{out_dir}" —— {e}', file=sys.stderr)
            print(f'   请确认是否有写入权限。', file=sys.stderr)
            sys.exit(1)

    # 5. 执行脱敏
    try:
        fmt, mapping = mask_document(
            args.input, args.output,
            strategy=args.strategy, seed=args.seed,
            skip_phone=args.skip_phone, map_path=args.map_file,
            as_text=args.as_text, ocr_lang=args.ocr_lang
        )
    except FileNotFoundError as e:
        print(f'错误：处理文件时找不到依赖或文件 —— {e}', file=sys.stderr)
        print(f'   请确认文件格式正确且未损坏。', file=sys.stderr)
        sys.exit(1)
    except ImportError as e:
        print(f'错误：缺少必要的 Python 库 —— {e}', file=sys.stderr)
        print(f'   请根据提示安装对应的库。常见：pip install python-docx python-pptx openpyxl PyPDF2 fpdf2', file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f'错误：脱敏过程中发生异常 —— {type(e).__name__}: {e}', file=sys.stderr)
        print(f'   如果问题持续，请检查输入文件是否损坏或格式是否正确。', file=sys.stderr)
        sys.exit(1)

    unique_count = len(mapping)
    print(f'✅ 脱敏完成: {unique_count} 个唯一值被替换 | 格式: {fmt} | 策略: {args.strategy}',
          file=sys.stderr)
    print(f'  输出文件: {args.output}', file=sys.stderr)
    if args.map_file:
        print(f'  映射表: {args.map_file}', file=sys.stderr)


if __name__ == '__main__':
    main()
