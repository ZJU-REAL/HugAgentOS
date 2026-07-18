"""File parsing utilities for chat attachments.

Supported formats:
  - PDF  : calls external file-parser API service
  - DOCX : pandoc (docx → markdown)
  - DOC / WPS : LibreOffice headless → DOCX, then pandoc
  - TXT  : UTF-8 / GBK direct decode
  - XLSX : openpyxl → markdown tables
  - XLS  : xlrd → markdown tables (fallback: LibreOffice → XLSX → openpyxl)
  - CSV  : csv module → markdown table

All public functions return a markdown/plain-text string, or raise RuntimeError
with a human-readable Chinese error message on failure.
"""

from __future__ import annotations

import csv
import io
import os
import subprocess
import tempfile
from typing import Optional

import requests


# ── config helpers (DB-first, env fallback) ──────────────────────────────────

def _svc_get(key: str, default: str = "") -> str:
    """Read from SystemConfigService (DB-first), fall back to env."""
    try:
        from core.services.system_config import SystemConfigService
        val = SystemConfigService.get_instance().get(key)
        if val is not None:
            return val.strip()
    except Exception:
        pass
    return (os.getenv(_ENV_MAP.get(key, ""), default) or default).strip()

# config_key → env var fallback mapping
_ENV_MAP = {
    "file_parser.api_url": "FILE_PARSER_API_URL",
    "file_parser.timeout": "FILE_PARSER_TIMEOUT",
    "file_parser.lang_list": "FILE_PARSER_LANG_LIST",
    "file_parser.backend": "FILE_PARSER_BACKEND",
    "file_parser.parse_method": "FILE_PARSER_PARSE_METHOD",
    "file_parser.formula_enable": "FILE_PARSER_FORMULA_ENABLE",
    "file_parser.table_enable": "FILE_PARSER_TABLE_ENABLE",
}


def _cfg_api_url() -> str:
    return _svc_get("file_parser.api_url")


def _cfg_timeout() -> int:
    try:
        return int(_svc_get("file_parser.timeout", "60"))
    except ValueError:
        return 60


def _cfg_parse_params() -> dict:
    return {
        "lang_list": _svc_get("file_parser.lang_list", "ch"),
        "backend": _svc_get("file_parser.backend", "pipeline"),
        "parse_method": _svc_get("file_parser.parse_method", "auto"),
        "formula_enable": _svc_get("file_parser.formula_enable", "true"),
        "table_enable": _svc_get("file_parser.table_enable", "true"),
    }


# ── PDF ───────────────────────────────────────────────────────────────────────

def parse_pdf(file_bytes: bytes, filename: str = "file.pdf") -> str:
    """Parse PDF via external file-parser API. Returns markdown text."""
    api_url = _cfg_api_url()
    if not api_url:
        raise RuntimeError("FILE_PARSER_API_URL 未配置，无法解析 PDF 文件")

    timeout = _cfg_timeout()

    try:
        resp = requests.post(
            api_url,
            files={"files": (filename, file_bytes, "application/pdf")},
            data=_cfg_parse_params(),
            timeout=timeout,
        )
        resp.raise_for_status()
    except requests.exceptions.Timeout:
        raise RuntimeError(f"PDF 解析服务超时（{timeout}s），请稍后重试")
    except requests.exceptions.RequestException as e:
        raise RuntimeError(f"PDF 解析服务请求失败: {e}")

    result = resp.json()
    results = result.get("results", {})
    if not results:
        raise RuntimeError("PDF 解析服务返回结果为空")

    title = next(iter(results))
    content = results[title].get("md_content", "")
    if not content:
        raise RuntimeError("PDF 解析服务返回内容为空")
    return content


# ── DOCX (via pandoc) ─────────────────────────────────────────────────────────

def _docx_bytes_to_markdown(docx_bytes: bytes) -> str:
    """Convert DOCX bytes to markdown string via pandoc."""
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.docx")
        with open(input_path, "wb") as f:
            f.write(docx_bytes)

        try:
            result = subprocess.run(
                ["pandoc", input_path, "-t", "markdown", "--wrap=none"],
                capture_output=True,
                text=True,
                timeout=60,
            )
        except FileNotFoundError:
            raise RuntimeError("pandoc 未安装，无法解析 Word 文档")
        except subprocess.TimeoutExpired:
            raise RuntimeError("pandoc 转换超时")

        if result.returncode != 0:
            raise RuntimeError(f"pandoc 转换失败: {result.stderr[:300]}")

        return result.stdout


def parse_docx(file_bytes: bytes) -> str:
    """Parse DOCX bytes → markdown via pandoc."""
    return _docx_bytes_to_markdown(file_bytes)


# ── DOC / WPS (LibreOffice → DOCX → pandoc) ──────────────────────────────────

def _convert_to_docx_bytes(file_bytes: bytes, suffix: str) -> bytes:
    """Use LibreOffice headless to convert DOC/WPS bytes → DOCX bytes."""
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, f"input{suffix}")
        with open(input_path, "wb") as f:
            f.write(file_bytes)

        try:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "docx",
                    "--outdir", tmpdir,
                    input_path,
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
        except FileNotFoundError:
            raise RuntimeError("LibreOffice 未安装，无法解析 DOC/WPS 文件")
        except subprocess.TimeoutExpired:
            raise RuntimeError("LibreOffice 转换超时")

        output_path = os.path.join(tmpdir, "input.docx")
        if not os.path.exists(output_path):
            raise RuntimeError(f"LibreOffice 转换失败: {result.stderr[:300]}")

        with open(output_path, "rb") as f:
            return f.read()


def parse_doc_wps(file_bytes: bytes, suffix: str) -> str:
    """Parse DOC/WPS via LibreOffice → DOCX → pandoc markdown."""
    docx_bytes = _convert_to_docx_bytes(file_bytes, suffix)
    return _docx_bytes_to_markdown(docx_bytes)


# ── TXT ───────────────────────────────────────────────────────────────────────

def parse_txt(file_bytes: bytes) -> str:
    """Decode text file bytes, trying UTF-8 then GBK then latin-1."""
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _deduplicate_headers(header_raw: list[str]) -> list[str]:
    """Disambiguate empty/duplicate header names so {placeholder} rendering
    downstream is unambiguous (e.g., two ``""`` columns become ``col``/``col_1``).
    """
    seen: dict[str, int] = {}
    cleaned: list[str] = []
    for h in header_raw:
        base = h if h else "col"
        if base in seen:
            seen[base] += 1
            cleaned.append(f"{base}_{seen[base]}")
        else:
            seen[base] = 0
            cleaned.append(base)
    return cleaned


def _rows_to_markdown(header: list[str], data_rows: list[list[str]],
                      sheet_title: str | None = None) -> str:
    """Convert header + data rows into a markdown table string."""
    lines: list[str] = []
    if sheet_title:
        lines.append(f"## Sheet: {sheet_title}")
        lines.append("")
    lines.append("| " + " | ".join(header) + " |")
    lines.append("|" + "|".join("---" for _ in header) + "|")
    for row in data_rows:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _open_xlsx(file_bytes: bytes):
    """Open a workbook in read-only mode; raise RuntimeError with a Chinese
    install hint if openpyxl is missing."""
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl 未安装，无法解析 XLSX 文件（pip install openpyxl）")
    return openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)


def _xlsx_sheet_to_markdown(sheet) -> str | None:
    """Render one openpyxl sheet to a markdown table, or None if empty.
    Returns None (not raise) so callers can decide whether emptiness is an
    error (single-sheet mode) or just a skip (whole-workbook mode)."""
    rows = list(sheet.iter_rows(values_only=True))
    if not rows:
        return None
    header = [str(c or "") for c in rows[0]]
    data_rows = [
        [str(c or "") for c in row]
        for row in rows[1:]
        if not all(c is None for c in row)
    ]
    return _rows_to_markdown(header, data_rows, sheet.title)


def parse_xlsx(file_bytes: bytes) -> str:
    """Parse XLSX bytes → markdown tables (one per sheet)."""
    wb = _open_xlsx(file_bytes)
    try:
        sections = [md for sheet in wb.worksheets if (md := _xlsx_sheet_to_markdown(sheet))]
    finally:
        wb.close()
    if not sections:
        raise RuntimeError("XLSX 文件无有效数据")
    return "\n\n".join(sections)


def parse_xlsx_sheet_names(file_bytes: bytes) -> list[str]:
    """Return the list of sheet names in workbook order. Cheap probe: does
    not iterate any rows, just reads workbook structure."""
    wb = _open_xlsx(file_bytes)
    try:
        return list(wb.sheetnames)
    finally:
        wb.close()


def parse_xlsx_single_sheet(file_bytes: bytes, sheet_name: str) -> str:
    """Parse only the named sheet → markdown table. Raises RuntimeError if
    the sheet does not exist or is empty."""
    wb = _open_xlsx(file_bytes)
    try:
        if sheet_name not in wb.sheetnames:
            raise RuntimeError(
                f"sheet 不存在: '{sheet_name}'。可选 sheet: {wb.sheetnames}"
            )
        md = _xlsx_sheet_to_markdown(wb[sheet_name])
        if md is None:
            raise RuntimeError(f"sheet '{sheet_name}' 无数据")
        return md
    finally:
        wb.close()


def _default_xlsx_max_rows() -> int:
    """Batch-execution row cap for parse_xlsx_structured.

    Read from BATCH_MAX_XLSX_ROWS env (default 2000). Each row is a
    separate per-row LLM call downstream, so the cap controls plan size
    and total runtime, not per-call context.
    """
    try:
        return max(1, int(os.getenv("BATCH_MAX_XLSX_ROWS", "2000")))
    except ValueError:
        return 2000


def _default_xlsx_max_cell_chars() -> int:
    """Per-cell character cap shared by parse_xlsx_structured and
    parse_xlsx_preview.

    Read from BATCH_MAX_XLSX_CELL_CHARS env (default 2000). Cells longer
    than this are truncated with an ellipsis. Larger values let long
    fields (summary/body/etc.) flow into per-row LLM prompts at the cost of
    more tokens per call and a larger plan.items JSONB row.
    """
    try:
        return max(50, int(os.getenv("BATCH_MAX_XLSX_CELL_CHARS", "2000")))
    except ValueError:
        return 2000


def parse_xlsx_structured(
    file_bytes: bytes,
    *,
    max_rows: Optional[int] = None,
    max_columns: int = 40,
    max_cell_chars: Optional[int] = None,
) -> tuple[list[str], list[dict], dict]:
    """Parse XLSX → (column_names, list_of_row_dicts, meta).

    Used for batch execution where each row becomes one task. Three
    independent safeguards bound the total plan size — but note each
    row is a separate downstream LLM call, so this is plan-runtime
    bounding, not per-call context bounding:

      - ``max_rows``       : drop rows beyond this index
                             (None → BATCH_MAX_XLSX_ROWS env, default 2000)
      - ``max_columns``    : drop columns beyond this index (default 40)
      - ``max_cell_chars`` : truncate any single cell value
                             (None → BATCH_MAX_XLSX_CELL_CHARS env, default 2000)

    Returns:
        (header_list, data_rows, meta) where ``meta`` describes what got
        truncated so the caller can surface a warning to the UI:
        ``{"total_rows", "total_columns", "truncated_rows": int,
        "truncated_columns": int, "truncated_cells": int}``.

    Reads only the first sheet (multi-sheet batch is out of scope).
    """
    if max_rows is None:
        max_rows = _default_xlsx_max_rows()
    if max_cell_chars is None:
        max_cell_chars = _default_xlsx_max_cell_chars()

    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl 未安装，无法解析 XLSX 文件（pip install openpyxl）")

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    try:
        sheet = wb.worksheets[0] if wb.worksheets else None
        if sheet is None:
            raise RuntimeError("XLSX 文件无 sheet")

        rows_iter = sheet.iter_rows(values_only=True)
        try:
            header_row = next(rows_iter)
        except StopIteration:
            raise RuntimeError("XLSX 文件无数据行") from None

        total_columns = len(header_row)
        truncated_columns = max(0, total_columns - max_columns)
        usable_cols = min(total_columns, max_columns)

        header = _deduplicate_headers([
            str(c).strip() if c is not None else f"col{idx + 1}"
            for idx, c in enumerate(header_row[:usable_cols])
        ])

        data: list[dict] = []
        truncated_cells = 0
        total_rows_seen = 0
        for row in rows_iter:
            total_rows_seen += 1
            if all(c is None for c in row):
                continue
            if len(data) >= max_rows:
                continue  # keep counting to report total_rows accurately
            d: dict = {}
            for idx, val in enumerate(row[:usable_cols]):
                if val is None:
                    d[header[idx]] = ""
                    continue
                s = str(val)
                if len(s) > max_cell_chars:
                    s = s[:max_cell_chars] + "…(已截断)"
                    truncated_cells += 1
                d[header[idx]] = s
            if any(v for v in d.values()):
                data.append(d)

        meta = {
            "total_rows": total_rows_seen,
            "total_columns": total_columns,
            "kept_rows": len(data),
            "kept_columns": usable_cols,
            "truncated_rows": max(0, total_rows_seen - max_rows),
            "truncated_columns": truncated_columns,
            "truncated_cells": truncated_cells,
            "max_cell_chars": max_cell_chars,
        }
        return header, data, meta
    finally:
        wb.close()


def parse_xlsx_preview(
    file_bytes: bytes,
    *,
    char_budget: int = 4000,
    max_columns: int = 40,
    max_cell_chars: Optional[int] = None,
) -> dict:
    """Parse XLSX into a bounded markdown preview + accurate total counts.

    Produces what the main agent's file_context hook should inject for an
    uploaded spreadsheet: a markdown table containing as many rows as fit
    within ``char_budget``, plus the *real* total_rows / total_columns so
    the agent knows the full size without seeing every row.

    The total_rows count walks the entire sheet (cheap — read_only mode
    streams the first column only). The preview stops as soon as the
    cumulative markdown length exceeds ``char_budget``.

    Returns:
        ``{
            "header": list[str],          # column names (post-dedup)
            "preview_md": str,            # markdown table, ≤ char_budget chars
            "total_rows": int,            # all data rows in sheet 1
            "total_columns": int,         # raw column count (before max_columns)
            "preview_rows": int,          # rows actually included in preview_md
            "truncated_columns": int,     # cols dropped to fit max_columns
            "truncated_cells": int,       # cells truncated to max_cell_chars
        }``

    Reads only the first sheet (matches parse_xlsx_structured).
    """
    if max_cell_chars is None:
        max_cell_chars = _default_xlsx_max_cell_chars()

    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl 未安装，无法解析 XLSX 文件（pip install openpyxl）")

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    try:
        sheet = wb.worksheets[0] if wb.worksheets else None
        if sheet is None:
            raise RuntimeError("XLSX 文件无 sheet")

        rows_iter = sheet.iter_rows(values_only=True)
        try:
            header_row = next(rows_iter)
        except StopIteration:
            raise RuntimeError("XLSX 文件无数据行") from None

        total_columns = len(header_row)
        truncated_columns = max(0, total_columns - max_columns)
        usable_cols = min(total_columns, max_columns)

        header = _deduplicate_headers([
            str(c).strip() if c is not None else f"col{idx + 1}"
            for idx, c in enumerate(header_row[:usable_cols])
        ])

        # Markdown header + separator (always present, counts toward budget)
        md_lines: list[str] = [
            "| " + " | ".join(header) + " |",
            "|" + "|".join("---" for _ in header) + "|",
        ]
        running_chars = sum(len(ln) + 1 for ln in md_lines)  # +1 for newline

        preview_rows = 0
        truncated_cells = 0
        total_rows_seen = 0
        budget_reached = False

        for row in rows_iter:
            total_rows_seen += 1
            if all(c is None for c in row):
                continue

            if budget_reached:
                continue  # keep counting total_rows accurately

            cells: list[str] = []
            for val in row[:usable_cols]:
                if val is None:
                    cells.append("")
                    continue
                s = str(val).replace("|", "\\|").replace("\n", " ")
                if len(s) > max_cell_chars:
                    s = s[:max_cell_chars] + "…"
                    truncated_cells += 1
                cells.append(s)

            # Pad short rows so the markdown table stays well-formed
            while len(cells) < len(header):
                cells.append("")

            line = "| " + " | ".join(cells) + " |"
            line_cost = len(line) + 1
            if running_chars + line_cost > char_budget and preview_rows > 0:
                # We have at least one row already; respect the budget.
                budget_reached = True
                continue
            md_lines.append(line)
            running_chars += line_cost
            preview_rows += 1

        return {
            "header": header,
            "preview_md": "\n".join(md_lines),
            "total_rows": total_rows_seen,
            "total_columns": total_columns,
            "preview_rows": preview_rows,
            "truncated_columns": truncated_columns,
            "truncated_cells": truncated_cells,
        }
    finally:
        wb.close()


def parse_xls(file_bytes: bytes) -> str:
    """Parse legacy XLS bytes → markdown tables.

    Tries xlrd first; falls back to LibreOffice → XLSX → openpyxl.
    """
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=file_bytes)
        sections: list[str] = []
        for sheet in wb.sheets():
            if sheet.nrows == 0:
                continue
            header = [str(sheet.cell_value(0, c) or "") for c in range(sheet.ncols)]
            data_rows = []
            for r in range(1, sheet.nrows):
                row = [str(sheet.cell_value(r, c) or "") for c in range(sheet.ncols)]
                if all(v == "" for v in row):
                    continue
                data_rows.append(row)
            sections.append(_rows_to_markdown(header, data_rows, sheet.name))
        if not sections:
            raise RuntimeError("XLS 文件无有效数据")
        return "\n\n".join(sections)
    except ImportError:
        pass  # xlrd not installed, try LibreOffice fallback

    # Fallback: LibreOffice → XLSX → openpyxl
    xlsx_bytes = _convert_xls_to_xlsx(file_bytes)
    return parse_xlsx(xlsx_bytes)


def _convert_xls_to_xlsx(file_bytes: bytes) -> bytes:
    """Use LibreOffice headless to convert XLS → XLSX."""
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.xls")
        with open(input_path, "wb") as f:
            f.write(file_bytes)

        try:
            result = subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to", "xlsx",
                    "--outdir", tmpdir,
                    input_path,
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
        except FileNotFoundError:
            raise RuntimeError(
                "xlrd 和 LibreOffice 均不可用，无法解析 XLS 文件"
                "（pip install xlrd 或安装 LibreOffice）"
            )
        except subprocess.TimeoutExpired:
            raise RuntimeError("LibreOffice 转换 XLS → XLSX 超时")

        output_path = os.path.join(tmpdir, "input.xlsx")
        if not os.path.exists(output_path):
            raise RuntimeError(f"LibreOffice 转换 XLS 失败: {result.stderr[:300]}")

        with open(output_path, "rb") as f:
            return f.read()


# ── CSV ──────────────────────────────────────────────────────────────────────

def parse_csv(file_bytes: bytes) -> str:
    """Parse CSV bytes → markdown table."""
    # Decode with encoding detection
    text = parse_txt(file_bytes)

    # Sniff dialect
    try:
        sample = text[:8192]
        dialect = csv.Sniffer().sniff(sample)
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect)
    rows = list(reader)

    if not rows:
        raise RuntimeError("CSV 文件无有效数据")

    header = [c.strip() for c in rows[0]]
    data_rows = []
    for row in rows[1:]:
        if all(c.strip() == "" for c in row):
            continue
        # Pad or trim to match header length
        padded = [c.strip() for c in row]
        while len(padded) < len(header):
            padded.append("")
        data_rows.append(padded[:len(header)])

    return _rows_to_markdown(header, data_rows)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _open_pptx(file_bytes: bytes):
    """Open a Presentation from bytes; raise RuntimeError with a Chinese
    install hint if python-pptx is missing."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError(
            "python-pptx 未安装，无法解析 PPTX 文件（pip install python-pptx）"
        )
    return Presentation(io.BytesIO(file_bytes))


def _extract_slide_text(slide) -> str:
    # Iterate in author/Z order so the rendered output mirrors what a
    # reader sees (title → body → footers), since python-pptx exposes
    # shapes in that order on the slide.shapes collection.
    parts: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            text = shape.text_frame.text
        except Exception:
            continue
        text = (text or "").strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def parse_pptx(file_bytes: bytes) -> str:
    """Parse PPTX bytes → markdown with one ``## Slide N`` section per slide."""
    p = _open_pptx(file_bytes)
    sections: list[str] = []
    for idx, slide in enumerate(p.slides):
        body = _extract_slide_text(slide) or "(此页无文本)"
        sections.append(f"## Slide {idx + 1}\n\n{body}")
    if not sections:
        raise RuntimeError("PPTX 文件无幻灯片")
    return "\n\n".join(sections)


def parse_pptx_slide_count(file_bytes: bytes) -> int:
    """Cheap probe: return the slide count without extracting any text."""
    return len(_open_pptx(file_bytes).slides)


def parse_pptx_slide(file_bytes: bytes, slide_index: int) -> str:
    """Return only the named slide's text (0-indexed). Raises RuntimeError
    on out-of-range index."""
    slides = _open_pptx(file_bytes).slides
    total = len(slides)
    if total == 0:
        raise RuntimeError("PPTX 文件无幻灯片")
    if slide_index < 0 or slide_index >= total:
        raise RuntimeError(
            f"slide_index 越界: {slide_index}（共 {total} 页，可选 0..{total - 1}）"
        )
    return _extract_slide_text(slides[slide_index]) or "(此页无文本)"


# ── Dispatcher ────────────────────────────────────────────────────────────────

_EXT_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc_wps",
    ".wps": "doc_wps",
    ".txt": "txt",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".csv": "csv",
    ".pptx": "pptx",
}

# Text-like extensions that don't need a real parser — just utf-8 decode.
# Covers HTML/Markdown/code/config formats the agent commonly produces.
# Mapping value "raw_text" routes to ``parse_raw_text``.
_RAW_TEXT_EXTS = {
    ".html", ".htm",
    ".md", ".markdown",
    ".json", ".jsonl", ".ndjson",
    ".yaml", ".yml",
    ".toml", ".ini", ".cfg",
    ".xml", ".svg",
    ".log",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash",
    ".css", ".scss",
    ".tex", ".rst",
    ".sql",
}

SUPPORTED_EXTENSIONS = list(_EXT_MAP.keys()) + sorted(_RAW_TEXT_EXTS)


def parse_raw_text(file_bytes: bytes) -> str:
    """UTF-8 decode for plain-text formats (HTML/MD/code/etc.).

    No structural extraction—just the source bytes as a string. Replaces
    undecodable bytes so we always return something usable.
    """
    return file_bytes.decode("utf-8", errors="replace")


def parse_file(file_bytes: bytes, filename: str) -> Optional[str]:
    """
    Dispatch to the correct parser by file extension.

    Returns extracted text/markdown, or None if the format is unsupported.
    Raises RuntimeError with a Chinese message on parse failure.
    """
    suffix = os.path.splitext(filename.lower())[1]
    kind = _EXT_MAP.get(suffix)

    if kind == "pdf":
        return parse_pdf(file_bytes, filename)
    elif kind == "docx":
        return parse_docx(file_bytes)
    elif kind == "doc_wps":
        return parse_doc_wps(file_bytes, suffix)
    elif kind == "txt":
        return parse_txt(file_bytes)
    elif kind == "xlsx":
        return parse_xlsx(file_bytes)
    elif kind == "xls":
        return parse_xls(file_bytes)
    elif kind == "csv":
        return parse_csv(file_bytes)
    elif kind == "pptx":
        return parse_pptx(file_bytes)
    elif suffix in _RAW_TEXT_EXTS:
        return parse_raw_text(file_bytes)
    else:
        return None
