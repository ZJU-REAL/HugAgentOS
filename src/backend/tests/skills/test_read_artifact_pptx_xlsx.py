"""Tests for read_artifact's xlsx-by-sheet and pptx-by-slide enhancements.

Covers:
- file_parser helpers: parse_xlsx_sheet_names / parse_xlsx_single_sheet,
  parse_pptx / parse_pptx_slide_count / parse_pptx_slide
- read_artifact metadata dispatch: _is_xlsx_meta / _is_pptx_meta
- read_artifact tool: mutual exclusion (sheet_name + slide_index),
  param-vs-filetype validation, sheet_names + slide_count propagation.

Skips DB/storage paths via monkeypatch.
"""

from __future__ import annotations

import io
import json
from typing import Any, Dict
from unittest.mock import patch

import pytest

# ── unit: file_parser xlsx helpers ────────────────────────────────────────


def _make_two_sheet_xlsx() -> bytes:
    """Build a 2-sheet xlsx in memory: 'Alpha' + 'Beta'."""
    import openpyxl
    wb = openpyxl.Workbook()
    s1 = wb.active
    s1.title = "Alpha"
    s1.append(["name", "score"])
    s1.append(["a", 1])
    s1.append(["b", 2])
    s2 = wb.create_sheet("Beta")
    s2.append(["city", "pop"])
    s2.append(["x", 100])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_parse_xlsx_sheet_names_returns_workbook_order():
    from core.content.file_parser import parse_xlsx_sheet_names
    names = parse_xlsx_sheet_names(_make_two_sheet_xlsx())
    assert names == ["Alpha", "Beta"]


def test_parse_xlsx_single_sheet_returns_only_that_sheet():
    from core.content.file_parser import parse_xlsx_single_sheet
    md = parse_xlsx_single_sheet(_make_two_sheet_xlsx(), "Beta")
    assert "city" in md
    assert "pop" in md
    assert "100" in md
    # Must NOT contain Alpha sheet content
    assert "Alpha" not in md or "Sheet: Beta" in md  # sheet header is "## Sheet: Beta"
    assert "score" not in md


def test_parse_xlsx_single_sheet_missing_raises():
    from core.content.file_parser import parse_xlsx_single_sheet
    with pytest.raises(RuntimeError, match="sheet 不存在"):
        parse_xlsx_single_sheet(_make_two_sheet_xlsx(), "NonExistent")


# ── unit: file-type detection ─────────────────────────────────────────────


def test_is_xlsx_meta_detects_by_mime_and_extension():
    from core.llm.tools.read_artifact_tool import _is_xlsx_meta
    assert _is_xlsx_meta({"mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"})
    assert _is_xlsx_meta({"filename": "report.xlsx"})
    assert _is_xlsx_meta({"filename": "REPORT.XLSX"})
    assert not _is_xlsx_meta({"filename": "doc.pdf"})
    assert not _is_xlsx_meta({})


def test_is_pptx_meta_detects_by_mime_and_extension():
    from core.llm.tools.read_artifact_tool import _is_pptx_meta
    assert _is_pptx_meta({"mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})
    assert _is_pptx_meta({"filename": "deck.pptx"})
    assert _is_pptx_meta({"filename": "DECK.PPTX"})
    assert not _is_pptx_meta({"filename": "notes.docx"})


# ── integration: read_artifact dispatch ───────────────────────────────────


@pytest.fixture
def registered_tool():
    """Register read_artifact onto a fresh Toolkit and return the callable."""
    from core.llm.tool_collector import ToolCollector
    from core.llm.tools.read_artifact_tool import register_read_artifact

    # AgentScope 2.0: register_* writes into the ToolCollector (duck-compatible with register_tool_function);
    # the tool's original callable lives at AllowedFunctionTool._func.
    c = ToolCollector()
    register_read_artifact(c, user_id="u1")
    return c.get_tool("read_artifact")._func


async def _call(func, **kwargs) -> Dict[str, Any]:
    """Invoke the tool and return the parsed JSON dict from its TextBlock."""
    resp = await func(**kwargs)
    block = resp.content[0]
    text = block["text"] if isinstance(block, dict) else block.text
    return json.loads(text)


@pytest.mark.asyncio
async def test_mutual_exclusion_sheet_and_slide(registered_tool):
    out = await _call(registered_tool, file_id="ua_x", sheet_name="S1", slide_index=0)
    assert "error" in out
    assert "互斥" in out["error"]


@pytest.mark.asyncio
async def test_missing_artifact_returns_error(registered_tool):
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=None):
        out = await _call(registered_tool, file_id="ua_missing")
    assert "error" in out
    assert "不存在" in out["error"]


@pytest.mark.asyncio
async def test_xlsx_default_call_includes_sheet_names(registered_tool):
    """Without sheet_name, returns concatenated text + sheet_names list."""
    file_bytes = _make_two_sheet_xlsx()
    meta = {"filename": "two.xlsx", "mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta), \
         patch("core.content.artifact_reader.fetch_parsed_text", return_value="full concatenated md"), \
         patch("core.llm.hooks._download_artifact_bytes", return_value=file_bytes):
        out = await _call(registered_tool, file_id="ua_xlsx")
    assert out.get("sheet_names") == ["Alpha", "Beta"]
    assert out["content"] == "full concatenated md"


@pytest.mark.asyncio
async def test_xlsx_with_sheet_name_returns_only_that_sheet(registered_tool):
    file_bytes = _make_two_sheet_xlsx()
    meta = {"filename": "two.xlsx", "mime_type": ""}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta), \
         patch("core.llm.hooks._download_artifact_bytes", return_value=file_bytes):
        out = await _call(registered_tool, file_id="ua_xlsx", sheet_name="Beta", limit=20000)
    assert out.get("sheet_names") == ["Alpha", "Beta"]
    assert "city" in out["content"]
    assert "score" not in out["content"]


@pytest.mark.asyncio
async def test_xlsx_unknown_sheet_returns_error_with_names(registered_tool):
    file_bytes = _make_two_sheet_xlsx()
    meta = {"filename": "two.xlsx", "mime_type": ""}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta), \
         patch("core.llm.hooks._download_artifact_bytes", return_value=file_bytes):
        out = await _call(registered_tool, file_id="ua_xlsx", sheet_name="Gamma")
    assert "error" in out
    assert "Gamma" in out["error"]
    assert out.get("sheet_names") == ["Alpha", "Beta"]


def _make_three_slide_pptx() -> bytes:
    """Build a 3-slide pptx in memory with title text on each slide."""
    from pptx import Presentation
    p = Presentation()
    titles = ["Hello", "World", "Goodbye"]
    for t in titles:
        slide = p.slides.add_slide(p.slide_layouts[5])  # title-only layout
        slide.shapes.title.text = t
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


def test_parse_pptx_emits_one_section_per_slide():
    from core.content.file_parser import parse_pptx
    md = parse_pptx(_make_three_slide_pptx())
    assert md.count("## Slide ") == 3
    assert "Hello" in md and "World" in md and "Goodbye" in md


def test_parse_pptx_slide_count_returns_int():
    from core.content.file_parser import parse_pptx_slide_count
    assert parse_pptx_slide_count(_make_three_slide_pptx()) == 3


def test_parse_pptx_slide_extracts_single_slide():
    from core.content.file_parser import parse_pptx_slide
    text = parse_pptx_slide(_make_three_slide_pptx(), 1)
    assert "World" in text
    assert "Hello" not in text and "Goodbye" not in text


def test_parse_pptx_slide_out_of_range_raises():
    from core.content.file_parser import parse_pptx_slide
    with pytest.raises(RuntimeError, match="越界"):
        parse_pptx_slide(_make_three_slide_pptx(), 10)


@pytest.mark.asyncio
async def test_pptx_default_returns_full_deck_with_slide_count(registered_tool):
    """Cold path: empty parsed_text → fetch_parsed_text triggers parse_pptx
    via the artifact_reader pipeline. We mock that and assert slide_count
    is reported from the in-process probe."""
    file_bytes = _make_three_slide_pptx()
    meta = {"filename": "deck.pptx", "mime_type": ""}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta), \
         patch("core.content.artifact_reader.fetch_parsed_text", return_value="full deck text"), \
         patch("core.llm.hooks._download_artifact_bytes", return_value=file_bytes):
        out = await _call(registered_tool, file_id="ua_pptx", limit=20000)
    assert out["content"] == "full deck text"
    assert out["slide_count"] == 3


@pytest.mark.asyncio
async def test_pptx_with_slide_index_extracts_only_that_slide(registered_tool):
    file_bytes = _make_three_slide_pptx()
    meta = {"filename": "deck.pptx", "mime_type": ""}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta), \
         patch("core.llm.hooks._download_artifact_bytes", return_value=file_bytes):
        out = await _call(registered_tool, file_id="ua_pptx", slide_index=1, limit=2000)
    assert "World" in out["content"]
    assert "Hello" not in out["content"]
    assert out["slide_count"] == 3
    assert out["slide_index"] == 1


@pytest.mark.asyncio
async def test_pptx_slide_index_out_of_range_returns_error_with_count(registered_tool):
    file_bytes = _make_three_slide_pptx()
    meta = {"filename": "deck.pptx", "mime_type": ""}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta), \
         patch("core.llm.hooks._download_artifact_bytes", return_value=file_bytes):
        out = await _call(registered_tool, file_id="ua_pptx", slide_index=99)
    assert "error" in out
    assert "越界" in out["error"]
    assert out["slide_count"] == 3


@pytest.mark.asyncio
async def test_non_office_with_sheet_or_slide_param_errors(registered_tool):
    meta = {"filename": "doc.pdf", "mime_type": "application/pdf"}
    with patch("core.content.artifact_reader.load_artifact_meta", return_value=meta):
        out = await _call(registered_tool, file_id="ua_pdf", sheet_name="Sheet1")
    assert "error" in out
    assert "xlsx" in out["error"]
